import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.drawing.image import Image as XLImage
import io
import tempfile
import os
import re
from collections import defaultdict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Set page configuration
st.set_page_config(
    page_title="Sales Report Generator",
    page_icon="📊",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for better styling
st.markdown("""
<style>
    .main-header {
        font-size: 2.5rem;
        font-weight: 700;
        color: #1f77b4;
        margin-bottom: 0.5rem;
    }
    .sub-header {
        font-size: 1.1rem;
        color: #666;
        margin-bottom: 2rem;
    }
    .metric-card {
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
        padding: 1rem;
        border-radius: 10px;
        color: white;
    }
    .section-divider {
        margin: 2rem 0;
        border-top: 2px solid #e0e0e0;
    }
    .download-section {
        background-color: #f8f9fa;
        padding: 1.5rem;
        border-radius: 10px;
        margin: 1rem 0;
    }
    .stTabs [data-baseweb="tab-list"] {
        gap: 2rem;
    }
    .stTabs [data-baseweb="tab"] {
        padding: 0.5rem 1.5rem;
        font-weight: 600;
    }
</style>
""", unsafe_allow_html=True)

# Header
st.markdown('<p class="main-header">📊 Sales Report Generator</p>', unsafe_allow_html=True)
st.markdown('<p class="sub-header">Comprehensive sales analytics with MTD/YTD metrics, brand performance, and SKU insights</p>', unsafe_allow_html=True)

# Sidebar for configuration
with st.sidebar:
    st.header("⚙️ Configuration")
    st.markdown("---")
    
    target_month = st.slider("📅 Target Month", 1, 12, 10, help="Month for MTD calculations")
    target_year = st.number_input("📆 Target Year", 2020, 2030, 2025, step=1, help="Year for calculations")
    comparison_year = st.number_input("📆 Comparison Year", 2020, 2030, 2024, step=1, help="Previous year for comparison")
    
    st.markdown("---")
    st.markdown("### 📖 Instructions")
    st.markdown("""
    1. Upload your Sales Summary Excel file
    2. Upload your ItemMaster Excel file (required)
    3. (Optional) Upload Sales Details & Customer List
    4. Review key performance metrics
    5. Download detailed reports
    6. Analyze brand, SKU, and customer performance
    """)
    
    st.markdown("---")
    st.markdown("### 📊 Available Reports")
    st.markdown("""
    - **Sales Summary**: Overall performance
    - **Top 10 Brands**: MTD & YTD rankings
    - **Top 20 SKUs**: Product-level insights
    - **Customer Channel**: Top 10 by channel (optional)
    """)

# Month names for display
month_names = ["", "January", "February", "March", "April", "May", "June", 
               "July", "August", "September", "October", "November", "December"]
month_name = month_names[target_month]

# File uploader with better styling
st.markdown("### 📁 Upload Your Data")

# Primary files (Sales Summary and ItemMaster)
col_upload1, col_upload2 = st.columns(2)

with col_upload1:
    uploaded_file = st.file_uploader(
        "📊 Sales Summary File",
        type=['xls', 'xlsx'],
        help="Upload the sales summary Excel file (e.g., summary 1-10.xls)",
        key="sales_file"
    )

with col_upload2:
    brand_mapping_file = st.file_uploader(
        "📋 ItemMaster File (Required)",
        type=['xls', 'xlsx'],
        help="Upload the ItemMaster Excel file containing ItemId and Brand columns. Required for accurate brand identification.",
        key="brand_file"
    )

# Additional required files for Customer Channel Report
col_upload3, col_upload4 = st.columns(2)

with col_upload3:
    sales_details_file = st.file_uploader(
        "📋 Sales Details File (Required)",
        type=['xls', 'xlsx'],
        help="Upload the Sales Details Excel file (e.g., sales detail 2024-2025.xls). Contains detailed transaction data.",
        key="sales_details_file"
    )

with col_upload4:
    customer_list_file = st.file_uploader(
        "👥 Sales Customer List File (Required)",
        type=['xls', 'xlsx'],
        help="Upload the SalesCustomerList Excel file with Channel column. Maps customers to their sales channels.",
        key="customer_list_file"
    )

def normalize_brand_name(brand_name):
    """Normalize brand name for fuzzy matching - removes spaces, lowercases"""
    if not brand_name:
        return ''
    # Remove all spaces and convert to lowercase for comparison
    return re.sub(r'\s+', '', str(brand_name).lower().strip())

def title_case_brand(brand_name):
    """Convert brand name to Title Case (capitalize each word)"""
    if not brand_name:
        return ''
    return ' '.join(word.capitalize() for word in str(brand_name).strip().split())

def load_brand_mapping(brand_file):
    """Load brand mapping from ItemMaster Excel file (uploaded directly by user)"""
    brand_map = {}  # ItemID -> Brand name
    brand_normalized_map = {}  # Normalized brand -> Canonical brand name
    
    try:
        # Read the ItemMaster file directly (user uploads this file)
        df = pd.read_excel(brand_file)
        
        # Find ItemId and Brand columns (based on actual file structure)
        item_id_col = None
        brand_col = None
        
        for col in df.columns:
            col_lower = str(col).lower().strip()
            if col_lower in ['itemid', 'item_id', 'item id', 'item']:
                item_id_col = col
            elif col_lower == 'brand':
                brand_col = col
        
        if item_id_col is None or brand_col is None:
            st.warning(f"⚠️ Could not find ItemId or Brand columns in the uploaded file. Found columns: {list(df.columns)}")
            return brand_map, brand_normalized_map
        
        # Build mapping: ItemID -> Brand (with Title Case)
        # Also build normalized brand map for fuzzy matching
        for _, row in df.iterrows():
            item_id = str(row[item_id_col]).strip() if pd.notna(row[item_id_col]) else ''
            brand = str(row[brand_col]).strip() if pd.notna(row[brand_col]) else ''
            if item_id and brand:
                # Normalize and title case the brand name
                normalized = normalize_brand_name(brand)
                title_brand = title_case_brand(brand)
                
                # If we've seen this normalized brand before, use the canonical version
                if normalized in brand_normalized_map:
                    title_brand = brand_normalized_map[normalized]
                else:
                    brand_normalized_map[normalized] = title_brand
                
                brand_map[item_id.upper()] = title_brand
        
        st.success(f"✅ Loaded {len(brand_map)} brand mappings ({len(brand_normalized_map)} unique brands)")
        
    except Exception as e:
        st.warning(f"⚠️ Error loading brand mapping: {str(e)}")
    
    return brand_map, brand_normalized_map

def split_by_items(df):
    """Split dataframe by item numbers"""
    items_dict = {}
    current_item = None
    current_data = []
    column_names = df.columns

    def finalize_item(rows, item_key):
        if not rows:
            return None
        item_df = pd.DataFrame(rows, columns=column_names)
        item_df = item_df.dropna(axis=1, how='all')
        if 'Item Number' in item_df.columns:
            item_df['Item Number'] = item_df['Item Number'].ffill().bfill().fillna(item_key)
        else:
            item_df.insert(0, 'Item Number', item_key)
        return item_df
    
    for idx, row in df.iterrows():
        first_col_value = str(row.iloc[0]) if pd.notna(row.iloc[0]) else ''
        first_col_clean = first_col_value.strip()
        
        is_item_code = (first_col_clean and 
                       any(c.isalpha() for c in first_col_clean) and 
                       any(c.isdigit() for c in first_col_clean) and
                       'Item Total' not in first_col_value)
        
        if is_item_code:
            if current_item and current_data:
                item_df = finalize_item(current_data, current_item)
                if item_df is not None:
                    items_dict[current_item] = item_df
            
            current_item = first_col_clean
            current_data = [row.values]
        elif current_item:
            current_data.append(row.values)
            
            if 'Item Total:' in str(row.values):
                item_df = finalize_item(current_data, current_item)
                if item_df is not None:
                    items_dict[current_item] = item_df
                current_item = None
                current_data = []
    
    if current_item and current_data:
        item_df = finalize_item(current_data, current_item)
        if item_df is not None:
            items_dict[current_item] = item_df
    
    return items_dict

def calculate_sales_metrics(all_items, target_month, target_year):
    """Calculate MTD and YTD sales metrics"""
    mtd_sales = 0
    mtd_cost = 0
    ytd_sales = 0
    ytd_cost = 0
    items_with_data = 0
    
    for item_key, item_df in all_items.items():
        if item_df.empty:
            continue
        
        df = item_df.copy()
        
        # Find column names
        year_col = None
        period_col = None
        sales_col = None
        cost_col = None
        
        for col in df.columns:
            col_lower = str(col).strip().lower()
            if 'year' in col_lower and not year_col:
                year_col = col
            elif 'period' in col_lower and not period_col:
                period_col = col
            elif 'sales amount' in col_lower and not sales_col:
                sales_col = col
            elif 'cost of sales' in col_lower and not cost_col:
                cost_col = col
        
        if not year_col or not period_col or not sales_col:
            continue
        
        # Convert to numeric
        df[year_col] = pd.to_numeric(df[year_col], errors='coerce')
        df[period_col] = pd.to_numeric(df[period_col], errors='coerce')
        df[sales_col] = pd.to_numeric(df[sales_col], errors='coerce').fillna(0)
        if cost_col:
            df[cost_col] = pd.to_numeric(df[cost_col], errors='coerce').fillna(0)
        
        year_data = df[df[year_col] == target_year].copy()
        
        if not year_data.empty:
            items_with_data += 1
            
            # MTD
            mtd_data = year_data[year_data[period_col] == target_month]
            mtd_sales += mtd_data[sales_col].sum()
            if cost_col:
                mtd_cost += mtd_data[cost_col].sum()
            
            # YTD
            ytd_data = year_data[year_data[period_col] <= target_month]
            ytd_sales += ytd_data[sales_col].sum()
            if cost_col:
                ytd_cost += ytd_data[cost_col].sum()
    
    mtd_gp = ((mtd_sales - mtd_cost) / mtd_sales * 100) if mtd_sales > 0 else 0
    ytd_gp = ((ytd_sales - ytd_cost) / ytd_sales * 100) if ytd_sales > 0 else 0
    
    return {
        'MTD Gross Sales': mtd_sales,
        'MTD GP%': mtd_gp,
        'YTD Gross Sales': ytd_sales,
        'YTD GP%': ytd_gp,
        'items_processed': items_with_data
    }

def create_excel_report(results_current, results_previous, target_month, target_year, comparison_year, month_name):
    """Create Excel report with formatted table and charts"""
    
    # Calculate % achieved
    mtd_achieved = (results_current['MTD Gross Sales'] / results_previous['MTD Gross Sales'] * 100) if results_previous['MTD Gross Sales'] > 0 else 0
    ytd_achieved = (results_current['YTD Gross Sales'] / results_previous['YTD Gross Sales'] * 100) if results_previous['YTD Gross Sales'] > 0 else 0
    mtd_gp_achieved = "0%" if results_previous['MTD GP%'] == 0 else f"{(results_current['MTD GP%'] / results_previous['MTD GP%'] * 100):.0f}%"
    ytd_gp_achieved = "0%" if results_previous['YTD GP%'] == 0 else f"{(results_current['YTD GP%'] / results_previous['YTD GP%'] * 100):.0f}%"
    
    # Create workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales Summary Report"
    
    # Title
    ws['A1'] = f"Sales Summary Report - {month_name} {target_year}"
    ws['A1'].font = Font(size=14, bold=True)
    ws['A1'].alignment = Alignment(horizontal='center')
    ws.merge_cells('A1:E1')
    
    # Styling
    header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
    header_font = Font(color="FFFFFF", bold=True)
    row_current_fill = PatternFill(start_color="E7E6E6", end_color="E7E6E6", fill_type="solid")
    row_achieved_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
    row_budget_fill = PatternFill(start_color="D9E1F2", end_color="D9E1F2", fill_type="solid")
    
    border_style = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # Headers
    headers = ['Date', 'MTD Gross Sales', 'MTD GP%', 'YTD Gross Sales', 'YTD GP%']
    for col_num, header in enumerate(headers, 1):
        cell = ws.cell(row=3, column=col_num, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center')
        cell.border = border_style
    
    # Data rows
    data_rows = [
        [str(target_year), 
         results_current['MTD Gross Sales'], 
         results_current['MTD GP%'] / 100,
         results_current['YTD Gross Sales'], 
         results_current['YTD GP%'] / 100],
        [str(comparison_year), 
         results_previous['MTD Gross Sales'], 
         results_previous['MTD GP%'] / 100, 
         results_previous['YTD Gross Sales'], 
         results_previous['YTD GP%'] / 100],
        ['%Achieved', 
         mtd_achieved / 100 if results_previous['MTD Gross Sales'] > 0 else 0,
         float(mtd_gp_achieved.rstrip('%')) / 100 if mtd_gp_achieved != "0%" else 0,
         ytd_achieved / 100 if results_previous['YTD Gross Sales'] > 0 else 0,
         float(ytd_gp_achieved.rstrip('%')) / 100 if ytd_gp_achieved != "0%" else 0],
        [f'{target_year} Budget', '', '', '', ''],
        ['% Achieved', 0, 0, 0, 0]
    ]
    
    # Write data rows
    for row_idx, row_data in enumerate(data_rows, start=4):
        for col_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=row_idx, column=col_idx, value=value)
            cell.border = border_style
            cell.alignment = Alignment(horizontal='right' if col_idx > 1 else 'left', vertical='center')
            
            if row_data[0] == str(target_year):
                cell.fill = row_current_fill
            elif '%Achieved' in row_data[0]:
                cell.fill = row_achieved_fill
            elif 'Budget' in row_data[0]:
                cell.fill = row_budget_fill
            
            if col_idx > 1 and value != '':
                if col_idx in [2, 4]:
                    cell.number_format = '$#,##0.00'
                elif col_idx in [3, 5]:
                    cell.number_format = '0.00%'
    
    # Set column widths
    ws.column_dimensions['A'].width = 15
    ws.column_dimensions['B'].width = 18
    ws.column_dimensions['C'].width = 12
    ws.column_dimensions['D'].width = 18
    ws.column_dimensions['E'].width = 12
    ws.row_dimensions[1].height = 25
    ws.row_dimensions[3].height = 20
    
    # Create dashboard chart
    fig_dashboard, axes = plt.subplots(2, 2, figsize=(12, 9))
    fig_dashboard.suptitle(f'Sales Performance Dashboard - {month_name} {target_year}', fontsize=14, fontweight='bold')
    
    categories = [str(target_year), str(comparison_year)]
    mtd_sales_values = [results_current['MTD Gross Sales'], results_previous['MTD Gross Sales']]
    ytd_sales_values = [results_current['YTD Gross Sales'], results_previous['YTD Gross Sales']]
    mtd_gp_values = [results_current['MTD GP%'], results_previous['MTD GP%']]
    ytd_gp_values = [results_current['YTD GP%'], results_previous['YTD GP%']]
    colors = ['#2E86AB', '#A23B72']
    
    # MTD Gross Sales
    ax1 = axes[0, 0]
    bars1 = ax1.bar(categories, mtd_sales_values, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
    ax1.set_title('MTD Gross Sales', fontsize=11, fontweight='bold')
    ax1.set_ylabel('Sales Amount ($)', fontsize=9)
    ax1.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
    for bar in bars1:
        height = bar.get_height()
        ax1.text(bar.get_x() + bar.get_width()/2., height, f'${height:,.0f}',
                ha='center', va='bottom', fontsize=8, fontweight='bold')
    if results_previous['MTD Gross Sales'] > 0:
        ax1.text(0.5, max(mtd_sales_values) * 0.92, f'% Achieved: {mtd_achieved:.0f}%',
                ha='center', fontsize=9, bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
    ax1.grid(axis='y', alpha=0.3, linestyle='--')
    
    # YTD Gross Sales
    ax2 = axes[0, 1]
    bars2 = ax2.bar(categories, ytd_sales_values, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
    ax2.set_title('YTD Gross Sales', fontsize=11, fontweight='bold')
    ax2.set_ylabel('Sales Amount ($)', fontsize=9)
    ax2.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x:,.0f}'))
    for bar in bars2:
        height = bar.get_height()
        ax2.text(bar.get_x() + bar.get_width()/2., height, f'${height:,.0f}',
                ha='center', va='bottom', fontsize=8, fontweight='bold')
    if results_previous['YTD Gross Sales'] > 0:
        ax2.text(0.5, max(ytd_sales_values) * 0.92, f'% Achieved: {ytd_achieved:.0f}%',
                ha='center', fontsize=9, bbox=dict(boxstyle='round', facecolor='wheat', alpha=0.5))
    ax2.grid(axis='y', alpha=0.3, linestyle='--')
    
    # MTD GP%
    ax3 = axes[1, 0]
    bars3 = ax3.bar(categories, mtd_gp_values, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
    ax3.set_title('MTD Gross Profit %', fontsize=11, fontweight='bold')
    ax3.set_ylabel('GP %', fontsize=9)
    ax3.set_ylim(0, max(mtd_gp_values) * 1.2 if max(mtd_gp_values) > 0 else 100)
    for bar in bars3:
        height = bar.get_height()
        ax3.text(bar.get_x() + bar.get_width()/2., height, f'{height:.2f}%',
                ha='center', va='bottom', fontsize=8, fontweight='bold')
    if results_previous['MTD GP%'] > 0:
        gp_achieved = (results_current['MTD GP%'] / results_previous['MTD GP%'] * 100)
        ax3.text(0.5, max(mtd_gp_values) * 0.92, f'% Achieved: {gp_achieved:.0f}%',
                ha='center', fontsize=9, bbox=dict(boxstyle='round', facecolor='lightgreen', alpha=0.5))
    ax3.grid(axis='y', alpha=0.3, linestyle='--')
    
    # YTD GP%
    ax4 = axes[1, 1]
    bars4 = ax4.bar(categories, ytd_gp_values, color=colors, alpha=0.8, edgecolor='black', linewidth=1.5)
    ax4.set_title('YTD Gross Profit %', fontsize=11, fontweight='bold')
    ax4.set_ylabel('GP %', fontsize=9)
    ax4.set_ylim(0, max(ytd_gp_values) * 1.2 if max(ytd_gp_values) > 0 else 100)
    for bar in bars4:
        height = bar.get_height()
        ax4.text(bar.get_x() + bar.get_width()/2., height, f'{height:.2f}%',
                ha='center', va='bottom', fontsize=8, fontweight='bold')
    if results_previous['YTD GP%'] > 0:
        gp_achieved = (results_current['YTD GP%'] / results_previous['YTD GP%'] * 100)
        ax4.text(0.5, max(ytd_gp_values) * 0.92, f'% Achieved: {gp_achieved:.0f}%',
                ha='center', fontsize=9, bbox=dict(boxstyle='round', facecolor='lightgreen', alpha=0.5))
    ax4.grid(axis='y', alpha=0.3, linestyle='--')
    
    plt.tight_layout()
    
    # Save dashboard chart
    img_buffer1 = io.BytesIO()
    fig_dashboard.savefig(img_buffer1, format='png', dpi=150, bbox_inches='tight')
    img_buffer1.seek(0)
    plt.close(fig_dashboard)
    
    ws_chart1 = wb.create_sheet("Dashboard Chart")
    img1 = XLImage(img_buffer1)
    ws_chart1.add_image(img1, 'A1')
    
    # Create comparison chart
    fig_comparison, ax = plt.subplots(1, 1, figsize=(11, 6))
    x = np.arange(4)
    width = 0.35
    metrics = ['MTD Gross Sales\n(in thousands)', 'MTD GP%', 'YTD Gross Sales\n(in thousands)', 'YTD GP%']
    values_current = [
        results_current['MTD Gross Sales'] / 1000,
        results_current['MTD GP%'],
        results_current['YTD Gross Sales'] / 1000,
        results_current['YTD GP%']
    ]
    values_previous = [
        results_previous['MTD Gross Sales'] / 1000,
        results_previous['MTD GP%'],
        results_previous['YTD Gross Sales'] / 1000,
        results_previous['YTD GP%']
    ]
    
    bars1 = ax.bar(x - width/2, values_current, width, label=str(target_year), color='#2E86AB', alpha=0.8, edgecolor='black', linewidth=1.5)
    bars2 = ax.bar(x + width/2, values_previous, width, label=str(comparison_year), color='#A23B72', alpha=0.8, edgecolor='black', linewidth=1.5)
    
    ax.set_title('Year-over-Year Performance Comparison', fontsize=13, fontweight='bold')
    ax.set_ylabel('Value', fontsize=10)
    ax.set_xticks(x)
    ax.set_xticklabels(metrics, fontsize=9)
    ax.legend(fontsize=10, loc='upper left')
    ax.grid(axis='y', alpha=0.3, linestyle='--')
    
    for bars in [bars1, bars2]:
        for i, bar in enumerate(bars):
            height = bar.get_height()
            if i in [0, 2]:
                label = f'${height:,.0f}K'
            else:
                label = f'{height:.2f}%'
            ax.text(bar.get_x() + bar.get_width()/2., height,
                    label, ha='center', va='bottom', fontsize=7, fontweight='bold')
    
    plt.tight_layout()
    
    # Save comparison chart
    img_buffer2 = io.BytesIO()
    fig_comparison.savefig(img_buffer2, format='png', dpi=150, bbox_inches='tight')
    img_buffer2.seek(0)
    plt.close(fig_comparison)
    
    ws_chart2 = wb.create_sheet("Comparison Chart")
    img2 = XLImage(img_buffer2)
    ws_chart2.add_image(img2, 'A1')
    
    # Save to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    return output

def extract_brand_code(item_key):
    """Extract letter prefix from item code (e.g., AMO0002 -> AMO)"""
    code = str(item_key).split('_')[-1]
    m = re.match(r'([A-Za-z]+)', code)
    return m.group(1).upper() if m else code.upper()

def compute_item_year_metrics(item_df, year, month):
    """Compute metrics for a single item dataframe and a given year"""
    df = item_df.copy()
    year_col = None
    period_col = None
    sales_col = None
    cost_col = None
    desc_col = None
    
    for col in df.columns:
        col_lower = str(col).strip().lower()
        if 'year' in col_lower and not year_col:
            year_col = col
        elif 'period' in col_lower and not period_col:
            period_col = col
        elif 'sales amount' in col_lower and not sales_col:
            sales_col = col
        elif 'cost of sales' in col_lower and not cost_col:
            cost_col = col
        elif any(k in col_lower for k in ['description','item','name']) and not desc_col:
            desc_col = col

    if not year_col or not period_col or not sales_col:
        return {'mtd_sales':0.0,'mtd_cost':0.0,'ytd_sales':0.0,'ytd_cost':0.0,'desc':None}

    df[year_col] = pd.to_numeric(df[year_col], errors='coerce')
    df[period_col] = pd.to_numeric(df[period_col], errors='coerce')
    df[sales_col] = pd.to_numeric(df[sales_col], errors='coerce').fillna(0)
    if cost_col:
        df[cost_col] = pd.to_numeric(df[cost_col], errors='coerce').fillna(0)

    year_data = df[df[year_col] == year]
    if year_data.empty:
        return {'mtd_sales':0.0,'mtd_cost':0.0,'ytd_sales':0.0,'ytd_cost':0.0,'desc':None}

    mtd_data = year_data[year_data[period_col] == month]
    ytd_data = year_data[year_data[period_col] <= month]

    mtd_sales = float(mtd_data[sales_col].sum())
    mtd_cost = float(mtd_data[cost_col].sum()) if cost_col else 0.0
    ytd_sales = float(ytd_data[sales_col].sum())
    ytd_cost = float(ytd_data[cost_col].sum()) if cost_col else 0.0

    desc_val = None
    if desc_col is not None:
        non_nulls = df[desc_col].dropna().astype(str)
        if len(non_nulls):
            desc_val = non_nulls.iloc[0]

    return {'mtd_sales':mtd_sales,'mtd_cost':mtd_cost,'ytd_sales':ytd_sales,'ytd_cost':ytd_cost,'desc':desc_val}

def create_sku_report(all_items, target_month, target_year, comparison_year, month_name, report_type='MTD'):
    """Create Top-20 SKU MTD or YTD Performance report (excludes MX and MEI products)"""
    
    # Collect all SKUs with their metrics
    sku_metrics = []
    
    for item_key, item_df in all_items.items():
        if item_df.empty:
            continue

        sku_code, item_name = get_sku_code_and_name(item_df)
        if not sku_code:
            continue

        brand_code = extract_brand_code(sku_code)
        # Skip MX brand and items starting with MEI
        if brand_code == 'MX':
            continue
        if sku_code.upper().startswith('MEI'):
            continue
        
        # Compute metrics for both years
        metrics_current = compute_item_year_metrics(item_df, target_year, target_month)
        metrics_previous = compute_item_year_metrics(item_df, comparison_year, target_month)
        
        # Calculate GP%
        def compute_gp(sales, cost):
            return ((sales - cost) / sales * 100) if sales > 0 else 0.0
        
        current_mtd_gp = compute_gp(metrics_current['mtd_sales'], metrics_current['mtd_cost'])
        current_ytd_gp = compute_gp(metrics_current['ytd_sales'], metrics_current['ytd_cost'])
        previous_mtd_gp = compute_gp(metrics_previous['mtd_sales'], metrics_previous['mtd_cost'])
        previous_ytd_gp = compute_gp(metrics_previous['ytd_sales'], metrics_previous['ytd_cost'])
        
        # Store metrics
        sku_metrics.append({
            'code': sku_code,
            'name': item_name,
            f'{target_year}_mtd_sales': metrics_current['mtd_sales'],
            f'{target_year}_mtd_qty': 0,  # Quantity not available in current data structure
            f'{target_year}_mtd_gp': current_mtd_gp,
            f'{comparison_year}_mtd_sales': metrics_previous['mtd_sales'],
            f'{comparison_year}_mtd_qty': 0,
            f'{comparison_year}_mtd_gp': previous_mtd_gp,
            f'{target_year}_ytd_sales': metrics_current['ytd_sales'],
            f'{target_year}_ytd_qty': 0,
            f'{target_year}_ytd_gp': current_ytd_gp,
            f'{comparison_year}_ytd_sales': metrics_previous['ytd_sales'],
            f'{comparison_year}_ytd_qty': 0,
            f'{comparison_year}_ytd_gp': previous_ytd_gp,
        })
    
    # Convert to DataFrame
    sku_df = pd.DataFrame(sku_metrics)
    if sku_df.empty:
        return None
    
    # Sort by MTD or YTD sales and get top 20
    sort_key = f'{target_year}_{report_type.lower()}_sales'
    sku_df = sku_df.sort_values(by=sort_key, ascending=False).head(20)
    
    # Create Excel workbook
    wb = Workbook()
    ws = wb.active
    ws.title = f'Top 20 {report_type} SKUs'
    
    # Title
    ws.merge_cells('A1:I1')
    tcell = ws['A1']
    tcell.value = f'Top 20 {report_type} SKUs Performance'
    tcell.font = Font(size=14, bold=True)
    tcell.alignment = Alignment(horizontal='center')
    
    # Header row styling
    header_fill_yellow = PatternFill(start_color="FFFF00", end_color="FFFF00", fill_type="solid")
    header_fill_green = PatternFill(start_color="92D050", end_color="92D050", fill_type="solid")
    
    metric_lower = report_type.lower()
    
    headers = [
        ('', None),
        ('Code', None),
        ('Item Name', None),
        (f'{target_year} {report_type}\nGross Sales', header_fill_yellow),
        (f'{target_year} {report_type}\nQTY', header_fill_yellow),
        (f'{target_year} {report_type}\nGP%', header_fill_yellow),
        (f'{comparison_year} {report_type}\nGross Sales', header_fill_green),
        (f'{comparison_year} {report_type}\nQTY', header_fill_green),
        (f'{comparison_year} {report_type} GP%', header_fill_green)
    ]
    
    for c_idx, (h, fill) in enumerate(headers, start=1):
        cell = ws.cell(row=2, column=c_idx, value=h)
        cell.font = Font(bold=True, size=10)
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = Border(left=Side(style='thin'), right=Side(style='thin'), 
                           top=Side(style='thin'), bottom=Side(style='thin'))
        if fill:
            cell.fill = fill
    
    # Write data rows
    for r_idx, row in sku_df.reset_index(drop=True).iterrows():
        rank = r_idx + 1
        row_data = [
            rank,
            row['code'],
            row['name'],
            row[f'{target_year}_{metric_lower}_sales'],
            row[f'{target_year}_{metric_lower}_qty'],
            row[f'{target_year}_{metric_lower}_gp'] / 100,
            row[f'{comparison_year}_{metric_lower}_sales'],
            row[f'{comparison_year}_{metric_lower}_qty'],
            row[f'{comparison_year}_{metric_lower}_gp'] / 100
        ]
        
        excel_row = r_idx + 3
        
        for c_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=excel_row, column=c_idx, value=value)
            cell.border = Border(left=Side(style='thin'), right=Side(style='thin'), 
                               top=Side(style='thin'), bottom=Side(style='thin'))
            
            # Apply formatting
            if c_idx in [4, 7]:  # Sales columns
                cell.number_format = '#,##0.00'
                cell.alignment = Alignment(horizontal='right')
            elif c_idx in [5, 8]:  # QTY columns
                cell.number_format = '#,##0'
                cell.alignment = Alignment(horizontal='right')
            elif c_idx in [6, 9]:  # GP% columns
                cell.number_format = '0.00%'
                cell.alignment = Alignment(horizontal='right')
            elif c_idx == 1:  # Rank
                cell.alignment = Alignment(horizontal='center')
            elif c_idx == 2:  # Code
                cell.alignment = Alignment(horizontal='left')
            else:  # Item name
                cell.alignment = Alignment(horizontal='left')
    
    # Set column widths
    widths = [5, 12, 50, 15, 12, 12, 15, 12, 14]
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[ws.cell(row=2, column=i).column_letter].width = w
    
    ws.row_dimensions[1].height = 25
    ws.row_dimensions[2].height = 35
    
    # Create charts
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 8))
    fig.suptitle(f'Top 20 SKUs - {report_type} Performance ({month_name} {target_year})', 
                 fontsize=14, fontweight='bold')
    
    # Extract data for charts (top 20)
    skus_list = [f"{row['code'][:10]}..." if len(row['code']) > 10 else row['code'] 
                 for _, row in sku_df.head(20).iterrows()]
    sales_current = sku_df[f'{target_year}_{metric_lower}_sales'].head(20).tolist()
    sales_previous = sku_df[f'{comparison_year}_{metric_lower}_sales'].head(20).tolist()
    gp_current = sku_df[f'{target_year}_{metric_lower}_gp'].head(20).tolist()
    gp_previous = sku_df[f'{comparison_year}_{metric_lower}_gp'].head(20).tolist()
    
    # Chart 1: Sales Comparison
    y_pos = np.arange(len(skus_list))
    width = 0.35
    
    bars1 = ax1.barh(y_pos - width/2, sales_current, width, label=str(target_year), 
                    color='#2E86AB', alpha=0.8, edgecolor='black', linewidth=0.8)
    bars2 = ax1.barh(y_pos + width/2, sales_previous, width, label=str(comparison_year), 
                    color='#A23B72', alpha=0.8, edgecolor='black', linewidth=0.8)
    
    ax1.set_xlabel('Sales Amount ($)', fontsize=10, fontweight='bold')
    ax1.set_ylabel('SKU Code', fontsize=10, fontweight='bold')
    ax1.set_title(f'{report_type} Gross Sales', fontsize=11, fontweight='bold')
    ax1.set_yticks(y_pos)
    ax1.set_yticklabels(skus_list, fontsize=7)
    ax1.legend(fontsize=9)
    ax1.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x/1000:.0f}K'))
    ax1.grid(axis='x', alpha=0.3, linestyle='--')
    ax1.invert_yaxis()
    
    # Chart 2: GP% Comparison
    bars3 = ax2.barh(y_pos - width/2, gp_current, width, label=str(target_year), 
                    color='#2E86AB', alpha=0.8, edgecolor='black', linewidth=0.8)
    bars4 = ax2.barh(y_pos + width/2, gp_previous, width, label=str(comparison_year), 
                    color='#A23B72', alpha=0.8, edgecolor='black', linewidth=0.8)
    
    ax2.set_xlabel('GP %', fontsize=10, fontweight='bold')
    ax2.set_title(f'{report_type} Gross Profit %', fontsize=11, fontweight='bold')
    ax2.set_yticks(y_pos)
    ax2.set_yticklabels(skus_list, fontsize=7)
    ax2.legend(fontsize=9)
    ax2.grid(axis='x', alpha=0.3, linestyle='--')
    ax2.invert_yaxis()
    
    plt.tight_layout()
    
    # Save chart to bytes
    img_buffer = io.BytesIO()
    fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close(fig)
    
    # Add chart to new sheet
    ws_chart = wb.create_sheet(f"{report_type} Charts")
    img = XLImage(img_buffer)
    ws_chart.add_image(img, 'A1')
    
    # Save to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    return output, sku_df

def get_sku_code_and_name(item_df):
    """Extract the original SKU code and a readable item name from the dataframe."""
    if item_df.empty:
        return '', ''

    first_row = item_df.iloc[0]

    # Extract code from Item Number column
    code = ''
    for col in item_df.columns:
        col_lower = str(col).lower()
        if 'item' in col_lower and any(k in col_lower for k in ['no', 'number', '#', 'code']):
            val = first_row[col]
            if pd.notna(val):
                val_str = str(val).strip()
                if val_str:
                    code = val_str.split()[0]
                    break

    if not code:
        raw_value = str(first_row.iloc[0]).strip()
        match = re.search(r'([A-Za-z]+\d+)', raw_value)
        if match:
            code = match.group(1)
        elif raw_value:
            code = raw_value.split()[0]

    # Extract name - look through all columns for a descriptive text
    name = ''
    # First try specific description columns
    for col in item_df.columns:
        col_lower = str(col).lower()
        if any(key in col_lower for key in ['description', 'item name', 'item desc']):
            val = first_row[col]
            if pd.notna(val):
                val_str = str(val).strip()
                if val_str and (not code or not val_str.upper().startswith(code.upper())):
                    name = val_str
                    break
    
    # If no name found, search all columns for a string that looks like a description
    # (contains spaces and is longer than the code)
    if not name:
        for col in item_df.columns:
            if col == 'Item Number':
                continue
            val = first_row[col]
            if pd.notna(val):
                val_str = str(val).strip()
                # Check if this looks like a product description
                if val_str and len(val_str) > 10 and ' ' in val_str and not val_str.replace('.','').replace(',','').isdigit():
                    # Make sure it's not just the code repeated
                    if not code or not val_str.upper().startswith(code.upper()):
                        name = val_str
                        break
                    # If it starts with code, extract the remainder
                    elif code and val_str.upper().startswith(code.upper()):
                        remainder = val_str[len(code):].strip(' -:_')
                        if remainder and len(remainder) > 3:
                            name = remainder
                            break

    if not name:
        raw_value = str(first_row.iloc[0]).strip()
        if code and raw_value.upper().startswith(code.upper()):
            remainder = raw_value[len(code):].strip(' -:_')
            if remainder:
                name = remainder
        elif raw_value and raw_value != code:
            name = raw_value

    if not code:
        code = name or ''
    if not name:
        name = code

    return code.strip(), name.strip()

def create_brand_report(all_items, target_month, target_year, comparison_year, month_name, report_type='MTD', brand_mapping=None, brand_normalized_map=None):
    """Create Top-10 Brand MTD or YTD Performance report (excludes MX brand)
    
    Requires brand_mapping from ItemMaster sheet for accurate brand identification.
    Items without brand mapping are skipped.
    """
    
    if brand_mapping is None:
        brand_mapping = {}
    if brand_normalized_map is None:
        brand_normalized_map = {}
    
    # Group items by brand name from mapping
    # Use normalized names for grouping to handle variations like "SunRice", "Sun Rice", "Sunrice"
    brands = defaultdict(list)
    brand_display_names = {}  # normalized_name -> display_name (Title Case)
    skipped_items = 0
    
    for item_key, item_df in all_items.items():
        # Extract item code from item_key (format: "SheetName_ItemCode")
        item_code = str(item_key).split('_')[-1].upper()
        brand_code = extract_brand_code(item_key)
        
        # Skip MX brand from reports
        if brand_code == 'MX':
            continue
        
        # Look up brand name from mapping using item code
        brand_name = brand_mapping.get(item_code, None)
        
        # Skip items without brand mapping - no fallback to first 3 letters
        if not brand_name:
            skipped_items += 1
            continue
        
        # Normalize the brand name for grouping (handles SunRice, Sun Rice, Sunrice as same)
        normalized_brand = normalize_brand_name(brand_name)
        
        # Check if we have a canonical name for this normalized brand
        if normalized_brand in brand_normalized_map:
            display_name = brand_normalized_map[normalized_brand]
        elif normalized_brand in brand_display_names:
            display_name = brand_display_names[normalized_brand]
        else:
            display_name = title_case_brand(brand_name)
            brand_display_names[normalized_brand] = display_name
        
        # Group by normalized name
        brands[normalized_brand].append((item_key, item_df, brand_code, display_name))
    
    # Aggregate per brand
    brand_metrics = []
    for normalized_brand, items in brands.items():
        # Get the display name (first item's display name, they should all be the same)
        display_name = items[0][3] if items else normalized_brand
        
        # Get the most common brand_code for this brand (for reference)
        brand_codes = [bc for _, _, bc, _ in items]
        most_common_code = max(set(brand_codes), key=brand_codes.count) if brand_codes else display_name
        
        agg = {
            'brand_code': most_common_code,
            'brand_name': display_name,
            f'{target_year}_mtd_sales': 0.0,
            f'{target_year}_mtd_cost': 0.0,
            f'{target_year}_ytd_sales': 0.0,
            f'{target_year}_ytd_cost': 0.0,
            f'{comparison_year}_mtd_sales': 0.0,
            f'{comparison_year}_mtd_cost': 0.0,
            f'{comparison_year}_ytd_sales': 0.0,
            f'{comparison_year}_ytd_cost': 0.0,
        }
        
        for item_key, item_df, _, _ in items:
            metrics_current = compute_item_year_metrics(item_df, target_year, target_month)
            metrics_previous = compute_item_year_metrics(item_df, comparison_year, target_month)
            
            agg[f'{target_year}_mtd_sales'] += metrics_current['mtd_sales']
            agg[f'{target_year}_mtd_cost'] += metrics_current['mtd_cost']
            agg[f'{target_year}_ytd_sales'] += metrics_current['ytd_sales']
            agg[f'{target_year}_ytd_cost'] += metrics_current['ytd_cost']
            agg[f'{comparison_year}_mtd_sales'] += metrics_previous['mtd_sales']
            agg[f'{comparison_year}_mtd_cost'] += metrics_previous['mtd_cost']
            agg[f'{comparison_year}_ytd_sales'] += metrics_previous['ytd_sales']
            agg[f'{comparison_year}_ytd_cost'] += metrics_previous['ytd_cost']
        
        # Compute GP%
        def compute_gp(sales, cost):
            return ((sales - cost) / sales * 100) if sales > 0 else 0.0
        
        agg[f'{target_year}_mtd_gp'] = compute_gp(agg[f'{target_year}_mtd_sales'], agg[f'{target_year}_mtd_cost'])
        agg[f'{target_year}_ytd_gp'] = compute_gp(agg[f'{target_year}_ytd_sales'], agg[f'{target_year}_ytd_cost'])
        agg[f'{comparison_year}_mtd_gp'] = compute_gp(agg[f'{comparison_year}_mtd_sales'], agg[f'{comparison_year}_mtd_cost'])
        agg[f'{comparison_year}_ytd_gp'] = compute_gp(agg[f'{comparison_year}_ytd_sales'], agg[f'{comparison_year}_ytd_cost'])
        
        # % Achieved
        agg['mtd_achieved_pct'] = (agg[f'{target_year}_mtd_sales'] / agg[f'{comparison_year}_mtd_sales'] * 100) if agg[f'{comparison_year}_mtd_sales'] > 0 else 0.0
        agg['ytd_achieved_pct'] = (agg[f'{target_year}_ytd_sales'] / agg[f'{comparison_year}_ytd_sales'] * 100) if agg[f'{comparison_year}_ytd_sales'] > 0 else 0.0
        
        brand_metrics.append(agg)
    
    # Convert to DataFrame and get top 10
    bm_df = pd.DataFrame(brand_metrics)
    if bm_df.empty:
        return None
    
    # Sort by MTD or YTD based on report type
    sort_key = f'{target_year}_{report_type.lower()}_sales'
    bm_df = bm_df.sort_values(by=sort_key, ascending=False).head(10)
    
    # Create Excel workbook
    wb = Workbook()
    ws = wb.active
    ws.title = f'Top 10 {report_type} Brands'
    
    # Title
    ws.merge_cells('A1:I1')
    tcell = ws['A1']
    tcell.value = f'Top 10 {report_type} Brand Performance - {month_name} {target_year}'
    tcell.font = Font(size=14, bold=True)
    tcell.alignment = Alignment(horizontal='center')
    
    # Header row styling
    header_fill = PatternFill(start_color="FFEB9C", end_color="FFEB9C", fill_type="solid")
    green_fill = PatternFill(start_color="C6E0B4", end_color="C6E0B4", fill_type="solid")
    blue_fill = PatternFill(start_color="BDD7EE", end_color="BDD7EE", fill_type="solid")
    
    headers = ['', 'Brand', 
               f'{target_year} {report_type}\nGross Sales', f'{target_year} {report_type}\nGP%',
               f'{comparison_year} {report_type}\nGross Sales', f'{comparison_year} {report_type}\nGP%',
               f'{target_year} {report_type}\nBudget ($)', 
               f'% Achieved vs\n{comparison_year}', '% Achieved vs\nBudget']
    
    for c_idx, h in enumerate(headers, start=1):
        cell = ws.cell(row=2, column=c_idx, value=h)
        cell.font = Font(bold=True, size=10)
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = Border(left=Side(style='thin'), right=Side(style='thin'), 
                           top=Side(style='thin'), bottom=Side(style='thin'))
        
        # Color code headers
        if c_idx in [3, 4]:  # Current year columns
            cell.fill = header_fill
        elif c_idx in [5, 6]:  # Previous year columns
            cell.fill = green_fill
        elif c_idx in [7, 8, 9]:  # Budget/achieved columns
            cell.fill = blue_fill
    
    # Write data rows
    metric_lower = report_type.lower()
    for r_idx, row in bm_df.reset_index(drop=True).iterrows():
        rank = r_idx + 1
        brand_name = row['brand_name']
        curr_sales = row[f'{target_year}_{metric_lower}_sales']
        curr_gp = row[f'{target_year}_{metric_lower}_gp'] / 100
        prev_sales = row[f'{comparison_year}_{metric_lower}_sales']
        prev_gp = row[f'{comparison_year}_{metric_lower}_gp'] / 100
        achieved = row[f'{metric_lower}_achieved_pct'] / 100
        
        row_data = [rank, brand_name, curr_sales, curr_gp, prev_sales, prev_gp, None, achieved, None]
        
        r_idx = r_idx + 3  # Adjust for Excel row numbering (header is row 2, data starts at row 3)
        
        for c_idx, value in enumerate(row_data, start=1):
            cell = ws.cell(row=r_idx, column=c_idx, value=value)
            cell.border = Border(left=Side(style='thin'), right=Side(style='thin'), 
                               top=Side(style='thin'), bottom=Side(style='thin'))
            
            # Apply formatting
            if c_idx in [3, 5, 7]:  # Sales columns
                cell.number_format = '$#,##0.00'
                cell.alignment = Alignment(horizontal='right')
            elif c_idx in [4, 6, 8, 9]:  # Percentage columns
                cell.number_format = '0.00%'
                cell.alignment = Alignment(horizontal='right')
            elif c_idx == 1:  # Rank
                cell.alignment = Alignment(horizontal='center')
            else:  # Brand name
                cell.alignment = Alignment(horizontal='left')
    
    # Set column widths
    widths = [6, 22, 18, 12, 18, 12, 18, 14, 16]
    for i, w in enumerate(widths, start=1):
        ws.column_dimensions[ws.cell(row=2, column=i).column_letter].width = w
    
    ws.row_dimensions[1].height = 25
    ws.row_dimensions[2].height = 35
    
    # Create charts
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(14, 6))
    fig.suptitle(f'Top 10 Brands - {report_type} Performance ({month_name} {target_year})', 
                 fontsize=14, fontweight='bold')
    
    # Extract data for charts
    brands_list = bm_df['brand_name'].head(10).tolist()
    metric_lower = report_type.lower()
    sales_current = bm_df[f'{target_year}_{metric_lower}_sales'].head(10).tolist()
    sales_previous = bm_df[f'{comparison_year}_{metric_lower}_sales'].head(10).tolist()
    gp_current = bm_df[f'{target_year}_{metric_lower}_gp'].head(10).tolist()
    gp_previous = bm_df[f'{comparison_year}_{metric_lower}_gp'].head(10).tolist()
    
    # Chart 1: Sales Comparison
    y_pos = np.arange(len(brands_list))
    width = 0.35
    
    bars1 = ax1.barh(y_pos - width/2, sales_current, width, label=str(target_year), 
                    color='#2E86AB', alpha=0.8, edgecolor='black', linewidth=1)
    bars2 = ax1.barh(y_pos + width/2, sales_previous, width, label=str(comparison_year), 
                    color='#A23B72', alpha=0.8, edgecolor='black', linewidth=1)
    
    ax1.set_xlabel('Sales Amount ($)', fontsize=10, fontweight='bold')
    ax1.set_ylabel('Brand', fontsize=10, fontweight='bold')
    ax1.set_title(f'{report_type} Gross Sales', fontsize=11, fontweight='bold')
    ax1.set_yticks(y_pos)
    ax1.set_yticklabels(brands_list, fontsize=8)
    ax1.legend(fontsize=9)
    ax1.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'${x/1000:.0f}K'))
    ax1.grid(axis='x', alpha=0.3, linestyle='--')
    ax1.invert_yaxis()
    
    # Add value labels
    for bars in [bars1, bars2]:
        for bar in bars:
            width_val = bar.get_width()
            if width_val > 0:
                ax1.text(width_val, bar.get_y() + bar.get_height()/2.,
                        f'${width_val/1000:.0f}K',
                        ha='left', va='center', fontsize=7, fontweight='bold')
    
    # Chart 2: GP% Comparison
    bars3 = ax2.barh(y_pos - width/2, gp_current, width, label=str(target_year), 
                    color='#2E86AB', alpha=0.8, edgecolor='black', linewidth=1)
    bars4 = ax2.barh(y_pos + width/2, gp_previous, width, label=str(comparison_year), 
                    color='#A23B72', alpha=0.8, edgecolor='black', linewidth=1)
    
    ax2.set_xlabel('GP %', fontsize=10, fontweight='bold')
    ax2.set_title(f'{report_type} Gross Profit %', fontsize=11, fontweight='bold')
    ax2.set_yticks(y_pos)
    ax2.set_yticklabels(brands_list, fontsize=8)
    ax2.legend(fontsize=9)
    ax2.grid(axis='x', alpha=0.3, linestyle='--')
    ax2.invert_yaxis()
    
    # Add value labels
    for bars in [bars3, bars4]:
        for bar in bars:
            width_val = bar.get_width()
            if width_val > 0:
                ax2.text(width_val, bar.get_y() + bar.get_height()/2.,
                        f'{width_val:.1f}%',
                        ha='left', va='center', fontsize=7, fontweight='bold')
    
    plt.tight_layout()
    
    # Save chart to bytes
    img_buffer = io.BytesIO()
    fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close(fig)
    
    # Add chart to new sheet
    ws_chart = wb.create_sheet(f"{report_type} Charts")
    img = XLImage(img_buffer)
    ws_chart.add_image(img, 'A1')
    
    # Save to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    return output, bm_df, len(brands)

def create_powerpoint_presentation(results_current, results_previous, target_month, target_year, 
                                   comparison_year, month_name, brand_mtd_df, brand_ytd_df, 
                                   sku_mtd_df, sku_ytd_df, channel_results=None, sales_rep_results=None):
    """Create a comprehensive PowerPoint presentation with all sales insights"""
    
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"Sales Performance Report"
    title_p = title_frame.paragraphs[0]
    title_p.font.size = Pt(44)
    title_p.font.bold = True
    title_p.font.color.rgb = RGBColor(31, 119, 180)
    title_p.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"{month_name} {target_year}"
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.font.size = Pt(32)
    subtitle_p.font.color.rgb = RGBColor(100, 100, 100)
    subtitle_p.alignment = PP_ALIGN.CENTER
    
    # Add date
    date_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    from datetime import datetime
    date_frame.text = f"Generated on {datetime.now().strftime('%B %d, %Y')}"
    date_p = date_frame.paragraphs[0]
    date_p.font.size = Pt(14)
    date_p.font.color.rgb = RGBColor(150, 150, 150)
    date_p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary with Key Metrics
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Title only
    title2 = slide2.shapes.title
    title2.text = "Executive Summary"
    title2.text_frame.paragraphs[0].font.size = Pt(36)
    title2.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 119, 180)
    
    # Calculate metrics
    mtd_change = ((results_current['MTD Gross Sales'] / results_previous['MTD Gross Sales'] - 1) * 100) if results_previous['MTD Gross Sales'] > 0 else 0
    ytd_change = ((results_current['YTD Gross Sales'] / results_previous['YTD Gross Sales'] - 1) * 100) if results_previous['YTD Gross Sales'] > 0 else 0
    mtd_gp_change = results_current['MTD GP%'] - results_previous['MTD GP%']
    ytd_gp_change = results_current['YTD GP%'] - results_previous['YTD GP%']
    
    # Create metrics boxes
    metrics = [
        ("MTD Gross Sales", f"${results_current['MTD Gross Sales']:,.0f}", f"{mtd_change:+.1f}%"),
        ("MTD GP%", f"{results_current['MTD GP%']:.2f}%", f"{mtd_gp_change:+.2f}%"),
        ("YTD Gross Sales", f"${results_current['YTD Gross Sales']:,.0f}", f"{ytd_change:+.1f}%"),
        ("YTD GP%", f"{results_current['YTD GP%']:.2f}%", f"{ytd_gp_change:+.2f}%")
    ]
    
    left_start = 0.5
    top_start = 2
    box_width = 2.2
    box_height = 2
    gap = 0.15
    
    for i, (label, value, change) in enumerate(metrics):
        col = i % 2
        row = i // 2
        left = left_start + col * (box_width + gap)
        top = top_start + row * (box_height + gap)
        
        # Background box
        shape = slide2.shapes.add_shape(
            1,  # Rectangle
            Inches(left), Inches(top), Inches(box_width), Inches(box_height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 248, 255)
        shape.line.color.rgb = RGBColor(31, 119, 180)
        shape.line.width = Pt(2)
        
        # Label
        label_box = slide2.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.2), Inches(box_width - 0.2), Inches(0.5))
        label_frame = label_box.text_frame
        label_frame.text = label
        label_p = label_frame.paragraphs[0]
        label_p.font.size = Pt(14)
        label_p.font.bold = True
        label_p.font.color.rgb = RGBColor(50, 50, 50)
        label_p.alignment = PP_ALIGN.CENTER
        
        # Value
        value_box = slide2.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.8), Inches(box_width - 0.2), Inches(0.6))
        value_frame = value_box.text_frame
        value_frame.text = value
        value_p = value_frame.paragraphs[0]
        value_p.font.size = Pt(24)
        value_p.font.bold = True
        value_p.font.color.rgb = RGBColor(31, 119, 180)
        value_p.alignment = PP_ALIGN.CENTER
        
        # Change
        change_box = slide2.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.5), Inches(box_width - 0.2), Inches(0.4))
        change_frame = change_box.text_frame
        change_frame.text = change
        change_p = change_frame.paragraphs[0]
        change_p.font.size = Pt(16)
        change_p.font.bold = True
        # Color based on positive/negative
        if '+' in change:
            change_p.font.color.rgb = RGBColor(0, 128, 0)
        else:
            change_p.font.color.rgb = RGBColor(255, 0, 0)
        change_p.alignment = PP_ALIGN.CENTER
    
    # Slide 3: MTD vs YTD Comparison Chart
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    title3 = slide3.shapes.title
    title3.text = f"MTD vs YTD Performance - {target_year}"
    title3.text_frame.paragraphs[0].font.size = Pt(32)
    
    # Create comparison chart
    fig, axes = plt.subplots(1, 2, figsize=(10, 4))
    
    categories = [str(target_year), str(comparison_year)]
    x = np.arange(len(categories))
    width = 0.35
    
    # Sales chart
    mtd_sales = [results_current['MTD Gross Sales'], results_previous['MTD Gross Sales']]
    ytd_sales = [results_current['YTD Gross Sales'], results_previous['YTD Gross Sales']]
    
    axes[0].bar(x - width/2, [v/1000 for v in mtd_sales], width, label='MTD', color='#2E86AB', alpha=0.8)
    axes[0].bar(x + width/2, [v/1000 for v in ytd_sales], width, label='YTD', color='#A23B72', alpha=0.8)
    axes[0].set_ylabel('Sales (in thousands)', fontsize=10, fontweight='bold')
    axes[0].set_title('Gross Sales Comparison', fontsize=12, fontweight='bold')
    axes[0].set_xticks(x)
    axes[0].set_xticklabels(categories)
    axes[0].legend()
    axes[0].grid(axis='y', alpha=0.3)
    
    # GP% chart
    mtd_gp = [results_current['MTD GP%'], results_previous['MTD GP%']]
    ytd_gp = [results_current['YTD GP%'], results_previous['YTD GP%']]
    
    axes[1].bar(x - width/2, mtd_gp, width, label='MTD', color='#2E86AB', alpha=0.8)
    axes[1].bar(x + width/2, ytd_gp, width, label='YTD', color='#A23B72', alpha=0.8)
    axes[1].set_ylabel('GP%', fontsize=10, fontweight='bold')
    axes[1].set_title('Gross Profit % Comparison', fontsize=12, fontweight='bold')
    axes[1].set_xticks(x)
    axes[1].set_xticklabels(categories)
    axes[1].legend()
    axes[1].grid(axis='y', alpha=0.3)
    
    plt.tight_layout()
    
    # Save chart to image
    chart_buffer = io.BytesIO()
    fig.savefig(chart_buffer, format='png', dpi=150, bbox_inches='tight')
    chart_buffer.seek(0)
    plt.close(fig)
    
    # Add chart to slide
    slide3.shapes.add_picture(chart_buffer, Inches(0.5), Inches(2), width=Inches(9))
    
    # Slide 4: Top 10 Brands - MTD
    if brand_mtd_df is not None and len(brand_mtd_df) > 0:
        slide4 = prs.slides.add_slide(prs.slide_layouts[5])
        title4 = slide4.shapes.title
        title4.text = "Top 10 Brands - MTD Performance"
        title4.text_frame.paragraphs[0].font.size = Pt(32)
        
        # Create table
        rows = min(11, len(brand_mtd_df) + 1)  # Header + up to 10 rows
        cols = 5
        left = Inches(0.5)
        top = Inches(2)
        width = Inches(9)
        height = Inches(4.5)
        
        table = slide4.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Set column widths
        table.columns[0].width = Inches(0.8)
        table.columns[1].width = Inches(3)
        table.columns[2].width = Inches(2)
        table.columns[3].width = Inches(1.6)
        table.columns[4].width = Inches(1.6)
        
        # Header
        headers = ['Rank', 'Brand', f'{target_year} MTD Sales', f'{target_year} GP%', '% Achieved']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        # Data rows
        for idx, row in brand_mtd_df.reset_index(drop=True).head(10).iterrows():
            table.cell(idx + 1, 0).text = str(idx + 1)
            table.cell(idx + 1, 1).text = str(row['brand_name'])
            table.cell(idx + 1, 2).text = f"${row[f'{target_year}_mtd_sales']:,.0f}"
            table.cell(idx + 1, 3).text = f"{row[f'{target_year}_mtd_gp']:.2f}%"
            table.cell(idx + 1, 4).text = f"{row['mtd_achieved_pct']:.1f}%"
            
            for col_idx in range(5):
                cell = table.cell(idx + 1, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    # Slide 5: Top 10 Brands - YTD
    if brand_ytd_df is not None and len(brand_ytd_df) > 0:
        slide5 = prs.slides.add_slide(prs.slide_layouts[5])
        title5 = slide5.shapes.title
        title5.text = "Top 10 Brands - YTD Performance"
        title5.text_frame.paragraphs[0].font.size = Pt(32)
        
        rows = min(11, len(brand_ytd_df) + 1)
        cols = 5
        table = slide5.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4.5)).table
        
        table.columns[0].width = Inches(0.8)
        table.columns[1].width = Inches(3)
        table.columns[2].width = Inches(2)
        table.columns[3].width = Inches(1.6)
        table.columns[4].width = Inches(1.6)
        
        headers = ['Rank', 'Brand', f'{target_year} YTD Sales', f'{target_year} GP%', '% Achieved']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for idx, row in brand_ytd_df.reset_index(drop=True).head(10).iterrows():
            table.cell(idx + 1, 0).text = str(idx + 1)
            table.cell(idx + 1, 1).text = str(row['brand_name'])
            table.cell(idx + 1, 2).text = f"${row[f'{target_year}_ytd_sales']:,.0f}"
            table.cell(idx + 1, 3).text = f"{row[f'{target_year}_ytd_gp']:.2f}%"
            table.cell(idx + 1, 4).text = f"{row['ytd_achieved_pct']:.1f}%"
            
            for col_idx in range(5):
                cell = table.cell(idx + 1, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(10)
    
    # Slides 6-7: Top 20 SKUs - MTD (split across 2 slides, 10 each)
    if sku_mtd_df is not None and len(sku_mtd_df) > 0:
        # Slide 6: SKUs 1-10
        slide6 = prs.slides.add_slide(prs.slide_layouts[5])
        title6 = slide6.shapes.title
        title6.text = "Top 20 SKUs - MTD Performance (1-10)"
        title6.text_frame.paragraphs[0].font.size = Pt(32)
        
        rows = min(11, len(sku_mtd_df) + 1)  # First 10 SKUs
        cols = 5
        table = slide6.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4.5)).table
        
        table.columns[0].width = Inches(0.6)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(3)
        table.columns[3].width = Inches(2)
        table.columns[4].width = Inches(2)
        
        headers = ['Rank', 'Code', 'Item Name', f'{target_year} MTD Sales', f'{target_year} GP%']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for idx, row in sku_mtd_df.reset_index(drop=True).head(10).iterrows():
            table.cell(idx + 1, 0).text = str(idx + 1)
            table.cell(idx + 1, 1).text = str(row['code'])
            table.cell(idx + 1, 2).text = str(row['name'])[:35] + '...' if len(str(row['name'])) > 35 else str(row['name'])
            table.cell(idx + 1, 3).text = f"${row[f'{target_year}_mtd_sales']:,.0f}"
            table.cell(idx + 1, 4).text = f"{row[f'{target_year}_mtd_gp']:.2f}%"
            
            for col_idx in range(5):
                cell = table.cell(idx + 1, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Slide 7: SKUs 11-20 (if available)
        if len(sku_mtd_df) > 10:
            slide7 = prs.slides.add_slide(prs.slide_layouts[5])
            title7 = slide7.shapes.title
            title7.text = "Top 20 SKUs - MTD Performance (11-20)"
            title7.text_frame.paragraphs[0].font.size = Pt(32)
            
            remaining_rows = min(11, len(sku_mtd_df) - 10 + 1)
            table = slide7.shapes.add_table(remaining_rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4.5)).table
            
            table.columns[0].width = Inches(0.6)
            table.columns[1].width = Inches(1.2)
            table.columns[2].width = Inches(3)
            table.columns[3].width = Inches(2)
            table.columns[4].width = Inches(2)
            
            # Header
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # Data rows 11-20
            for idx, row in sku_mtd_df.reset_index(drop=True).iloc[10:20].iterrows():
                row_num = idx - 9  # Adjust for table position
                table.cell(row_num, 0).text = str(idx + 1)
                table.cell(row_num, 1).text = str(row['code'])
                table.cell(row_num, 2).text = str(row['name'])[:35] + '...' if len(str(row['name'])) > 35 else str(row['name'])
                table.cell(row_num, 3).text = f"${row[f'{target_year}_mtd_sales']:,.0f}"
                table.cell(row_num, 4).text = f"{row[f'{target_year}_mtd_gp']:.2f}%"
                
                for col_idx in range(5):
                    cell = table.cell(row_num, col_idx)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    # Slides 8-9: Top 20 SKUs - YTD (split across 2 slides, 10 each)
    if sku_ytd_df is not None and len(sku_ytd_df) > 0:
        # Slide 8: SKUs 1-10
        slide8 = prs.slides.add_slide(prs.slide_layouts[5])
        title8 = slide8.shapes.title
        title8.text = "Top 20 SKUs - YTD Performance (1-10)"
        title8.text_frame.paragraphs[0].font.size = Pt(32)
        
        rows = min(11, len(sku_ytd_df) + 1)
        cols = 5
        table = slide8.shapes.add_table(rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4.5)).table
        
        table.columns[0].width = Inches(0.6)
        table.columns[1].width = Inches(1.2)
        table.columns[2].width = Inches(3)
        table.columns[3].width = Inches(2)
        table.columns[4].width = Inches(2)
        
        headers = ['Rank', 'Code', 'Item Name', f'{target_year} YTD Sales', f'{target_year} GP%']
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.size = Pt(10)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
        
        for idx, row in sku_ytd_df.reset_index(drop=True).head(10).iterrows():
            table.cell(idx + 1, 0).text = str(idx + 1)
            table.cell(idx + 1, 1).text = str(row['code'])
            table.cell(idx + 1, 2).text = str(row['name'])[:35] + '...' if len(str(row['name'])) > 35 else str(row['name'])
            table.cell(idx + 1, 3).text = f"${row[f'{target_year}_ytd_sales']:,.0f}"
            table.cell(idx + 1, 4).text = f"{row[f'{target_year}_ytd_gp']:.2f}%"
            
            for col_idx in range(5):
                cell = table.cell(idx + 1, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
        
        # Slide 9: SKUs 11-20 (if available)
        if len(sku_ytd_df) > 10:
            slide9 = prs.slides.add_slide(prs.slide_layouts[5])
            title9 = slide9.shapes.title
            title9.text = "Top 20 SKUs - YTD Performance (11-20)"
            title9.text_frame.paragraphs[0].font.size = Pt(32)
            
            remaining_rows = min(11, len(sku_ytd_df) - 10 + 1)
            table = slide9.shapes.add_table(remaining_rows, cols, Inches(0.5), Inches(2), Inches(9), Inches(4.5)).table
            
            table.columns[0].width = Inches(0.6)
            table.columns[1].width = Inches(1.2)
            table.columns[2].width = Inches(3)
            table.columns[3].width = Inches(2)
            table.columns[4].width = Inches(2)
            
            # Header
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            # Data rows 11-20
            for idx, row in sku_ytd_df.reset_index(drop=True).iloc[10:20].iterrows():
                row_num = idx - 9  # Adjust for table position
                table.cell(row_num, 0).text = str(idx + 1)
                table.cell(row_num, 1).text = str(row['code'])
                table.cell(row_num, 2).text = str(row['name'])[:35] + '...' if len(str(row['name'])) > 35 else str(row['name'])
                table.cell(row_num, 3).text = f"${row[f'{target_year}_ytd_sales']:,.0f}"
                table.cell(row_num, 4).text = f"{row[f'{target_year}_ytd_gp']:.2f}%"
                
                for col_idx in range(5):
                    cell = table.cell(row_num, col_idx)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
    
    # Slide: Channel Sales Performance
    if channel_results is not None:
        prev_year = target_year - 1
        channel_col = channel_results['channel_col']
        channel_metrics = channel_results['channel_metrics']
        
        # Filter out Unknown channels
        valid_channels = channel_metrics[channel_metrics[channel_col] != 'Unknown']
        
        if len(valid_channels) > 0:
            slide_channel = prs.slides.add_slide(prs.slide_layouts[5])
            title_channel = slide_channel.shapes.title
            title_channel.text = "Channel Sales Performance"
            title_channel.text_frame.paragraphs[0].font.size = Pt(32)
            title_channel.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 119, 180)
            
            # Create table for channel metrics
            rows = min(len(valid_channels) + 2, 12)  # Header + data + total (max 12 rows)
            cols = 7
            left = Inches(0.3)
            top = Inches(1.8)
            width = Inches(9.4)
            height = Inches(5)
            
            table = slide_channel.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths
            table.columns[0].width = Inches(1.8)
            table.columns[1].width = Inches(0.9)
            table.columns[2].width = Inches(1.4)
            table.columns[3].width = Inches(1.4)
            table.columns[4].width = Inches(0.9)
            table.columns[5].width = Inches(1.5)
            table.columns[6].width = Inches(1.5)
            
            # Header
            headers = ['Channel', 'Customers', f'{target_year} MTD', f'{prev_year} MTD', 'MTD %', f'{target_year} YTD', f'{prev_year} YTD']
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Data rows
            for idx, (_, ch_row) in enumerate(valid_channels.head(9).iterrows()):
                row_num = idx + 1
                table.cell(row_num, 0).text = str(ch_row[channel_col])[:20]
                table.cell(row_num, 1).text = str(int(ch_row['Customer_Count']))
                table.cell(row_num, 2).text = f"${ch_row[f'{target_year}_MTD_Sales']:,.0f}"
                table.cell(row_num, 3).text = f"${ch_row[f'{prev_year}_MTD_Sales']:,.0f}"
                table.cell(row_num, 4).text = f"{ch_row['MTD_Achieved_%']:.1f}%"
                table.cell(row_num, 5).text = f"${ch_row[f'{target_year}_YTD_Sales']:,.0f}"
                table.cell(row_num, 6).text = f"${ch_row[f'{prev_year}_YTD_Sales']:,.0f}"
                
                for col_idx in range(7):
                    cell = table.cell(row_num, col_idx)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Grand Total row
            gt = channel_results['grand_total']
            total_row = min(len(valid_channels.head(9)) + 1, 10)
            
            table.cell(total_row, 0).text = "GRAND TOTAL"
            table.cell(total_row, 1).text = str(gt['Customer_Count'])
            table.cell(total_row, 2).text = f"${gt[f'{target_year}_MTD_Sales']:,.0f}"
            table.cell(total_row, 3).text = f"${gt[f'{prev_year}_MTD_Sales']:,.0f}"
            table.cell(total_row, 4).text = f"{gt['MTD_Achieved_%']:.1f}%"
            table.cell(total_row, 5).text = f"${gt[f'{target_year}_YTD_Sales']:,.0f}"
            table.cell(total_row, 6).text = f"${gt[f'{prev_year}_YTD_Sales']:,.0f}"
            
            for col_idx in range(7):
                cell = table.cell(total_row, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Slide: Sales Rep Performance
    if sales_rep_results is not None:
        prev_year = target_year - 1
        rep_col = sales_rep_results['sales_rep_col']
        rep_metrics = sales_rep_results['sales_rep_metrics']
        
        # Filter out Unassigned reps
        valid_reps = rep_metrics[rep_metrics[rep_col] != 'Unassigned']
        
        if len(valid_reps) > 0:
            slide_rep = prs.slides.add_slide(prs.slide_layouts[5])
            title_rep = slide_rep.shapes.title
            title_rep.text = "Sales Rep Performance"
            title_rep.text_frame.paragraphs[0].font.size = Pt(32)
            title_rep.text_frame.paragraphs[0].font.color.rgb = RGBColor(31, 119, 180)
            
            # Create table for sales rep metrics
            rows = min(len(valid_reps) + 2, 12)  # Header + data + total (max 12 rows)
            cols = 7
            left = Inches(0.3)
            top = Inches(1.8)
            width = Inches(9.4)
            height = Inches(5)
            
            table = slide_rep.shapes.add_table(rows, cols, left, top, width, height).table
            
            # Set column widths
            table.columns[0].width = Inches(1.8)
            table.columns[1].width = Inches(0.9)
            table.columns[2].width = Inches(1.4)
            table.columns[3].width = Inches(1.4)
            table.columns[4].width = Inches(0.9)
            table.columns[5].width = Inches(1.5)
            table.columns[6].width = Inches(1.5)
            
            # Header
            headers = ['Sales Rep', 'Customers', f'{target_year} MTD', f'{prev_year} MTD', 'MTD %', f'{target_year} YTD', f'{prev_year} YTD']
            for col_idx, header in enumerate(headers):
                cell = table.cell(0, col_idx)
                cell.text = header
                cell.text_frame.paragraphs[0].font.size = Pt(10)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(31, 119, 180)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Data rows
            for idx, (_, rep_row) in enumerate(valid_reps.head(9).iterrows()):
                row_num = idx + 1
                table.cell(row_num, 0).text = str(rep_row[rep_col])[:20]
                table.cell(row_num, 1).text = str(int(rep_row['Customer_Count']))
                table.cell(row_num, 2).text = f"${rep_row[f'{target_year}_MTD_Sales']:,.0f}"
                table.cell(row_num, 3).text = f"${rep_row[f'{prev_year}_MTD_Sales']:,.0f}"
                table.cell(row_num, 4).text = f"{rep_row['MTD_Achieved_%']:.1f}%"
                table.cell(row_num, 5).text = f"${rep_row[f'{target_year}_YTD_Sales']:,.0f}"
                table.cell(row_num, 6).text = f"${rep_row[f'{prev_year}_YTD_Sales']:,.0f}"
                
                for col_idx in range(7):
                    cell = table.cell(row_num, col_idx)
                    cell.text_frame.paragraphs[0].font.size = Pt(9)
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Grand Total row
            gt_rep = sales_rep_results['grand_total']
            total_row = min(len(valid_reps.head(9)) + 1, 10)
            
            table.cell(total_row, 0).text = "GRAND TOTAL"
            table.cell(total_row, 1).text = str(gt_rep['Customer_Count'])
            table.cell(total_row, 2).text = f"${gt_rep[f'{target_year}_MTD_Sales']:,.0f}"
            table.cell(total_row, 3).text = f"${gt_rep[f'{prev_year}_MTD_Sales']:,.0f}"
            table.cell(total_row, 4).text = f"{gt_rep['MTD_Achieved_%']:.1f}%"
            table.cell(total_row, 5).text = f"${gt_rep[f'{target_year}_YTD_Sales']:,.0f}"
            table.cell(total_row, 6).text = f"${gt_rep[f'{prev_year}_YTD_Sales']:,.0f}"
            
            for col_idx in range(7):
                cell = table.cell(total_row, col_idx)
                cell.text_frame.paragraphs[0].font.size = Pt(9)
                cell.text_frame.paragraphs[0].font.bold = True
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(221, 235, 247)
                cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Save presentation to bytes
    ppt_output = io.BytesIO()
    prs.save(ppt_output)
    ppt_output.seek(0)
    
    return ppt_output


def generate_top10_customers_by_channel(
    sales_file,
    customer_list_file,
    target_month: int,
    target_year: int
) -> dict:
    """
    Generate Top 10 Customer performance reports by Channel for both MTD and YTD.
    Also calculates MTD/YTD totals for each Channel.
    
    Args:
        sales_file: Uploaded sales details Excel file object
        customer_list_file: Uploaded SalesCustomerList Excel file object
        target_month: Month for MTD calculation
        target_year: Year for calculations
    
    Returns:
        Dictionary with channel-wise top 10 customer data for MTD and YTD,
        plus channel totals
    """
    prev_year = target_year - 1
    
    # ============================================================
    # STEP 1: Parse sales data
    # ============================================================
    df_raw = pd.read_excel(sales_file, header=None)
    
    col_period = 1
    col_sales_amount = 17
    
    data_rows = []
    current_customer = None
    current_customer_name = None
    
    for idx in range(11, len(df_raw)):
        row = df_raw.iloc[idx]
        first_cell = row.iloc[0]
        
        if pd.isna(first_cell):
            continue
        
        first_cell_str = str(first_cell).strip()
        
        if first_cell_str and not first_cell_str.isdigit() and not first_cell_str.startswith('Item'):
            if any(c.isalpha() for c in first_cell_str) and any(c.isdigit() for c in first_cell_str):
                current_customer = first_cell_str
                cust_name = row.iloc[7] if pd.notna(row.iloc[7]) else ''
                current_customer_name = str(cust_name).strip() if cust_name else current_customer
                continue
        
        if first_cell_str.isdigit() and len(first_cell_str) == 4:
            year = int(first_cell_str)
            
            period_cell = row.iloc[col_period]
            month = None
            if pd.notna(period_cell):
                try:
                    month = int(float(period_cell))
                except (ValueError, TypeError):
                    month = None
            
            sales_amount = row.iloc[col_sales_amount]
            
            if pd.notna(sales_amount) and current_customer and month is not None:
                try:
                    sales_value = float(sales_amount)
                    data_rows.append({
                        'Customer #': current_customer,
                        'Customer Name': current_customer_name,
                        'Year': year,
                        'Month': month,
                        'Sales Amount': sales_value
                    })
                except (ValueError, TypeError):
                    pass
    
    df_sales = pd.DataFrame(data_rows)
    
    if df_sales.empty:
        return None
    
    # ============================================================
    # STEP 2: Join with customer list to get Channel
    # ============================================================
    customer_df = pd.read_excel(customer_list_file)
    
    # Find Channel and Sales Rep columns
    channel_col = None
    sales_rep_col = None
    for col in customer_df.columns:
        if 'channel' in col.lower():
            channel_col = col
        if 'sales' in col.lower() and 'rep' in col.lower():
            sales_rep_col = col
    
    if not channel_col:
        st.warning("⚠️ Channel column not found in customer list!")
        return None
    
    # Create lookup for customer -> channel/sales rep
    customer_cols = ['Customer #', channel_col]
    if sales_rep_col:
        customer_cols.append(sales_rep_col)
    
    customer_lookup = customer_df[customer_cols].drop_duplicates(subset=['Customer #'])
    
    # Merge to add channel to sales data
    df_sales = df_sales.merge(customer_lookup, on='Customer #', how='left')
    
    # Fill missing channels with 'Unknown'
    df_sales[channel_col] = df_sales[channel_col].fillna('Unknown')
    
    # ============================================================
    # STEP 3: Calculate MTD and YTD metrics by customer
    # ============================================================
    # Filter for target years
    df_current = df_sales[df_sales['Year'] == target_year].copy()
    df_prev = df_sales[df_sales['Year'] == prev_year].copy()
    
    # Calculate MTD (specific month only)
    df_mtd_current = df_current[df_current['Month'] == target_month]
    df_mtd_prev = df_prev[df_prev['Month'] == target_month]
    
    # Calculate YTD (months 1 to target_month)
    df_ytd_current = df_current[df_current['Month'] <= target_month]
    df_ytd_prev = df_prev[df_prev['Month'] <= target_month]
    
    # Aggregate by customer
    def aggregate_by_customer(df, ch_col, sr_col):
        agg_cols = ['Customer #', 'Customer Name', ch_col]
        if sr_col and sr_col in df.columns:
            agg_cols.append(sr_col)
        
        result = df.groupby(agg_cols, dropna=False).agg({
            'Sales Amount': 'sum'
        }).reset_index()
        return result
    
    mtd_current_agg = aggregate_by_customer(df_mtd_current, channel_col, sales_rep_col)
    mtd_prev_agg = aggregate_by_customer(df_mtd_prev, channel_col, sales_rep_col)
    ytd_current_agg = aggregate_by_customer(df_ytd_current, channel_col, sales_rep_col)
    ytd_prev_agg = aggregate_by_customer(df_ytd_prev, channel_col, sales_rep_col)
    
    # Rename columns
    mtd_current_agg = mtd_current_agg.rename(columns={'Sales Amount': f'{target_year}_MTD_Sales'})
    mtd_prev_agg = mtd_prev_agg.rename(columns={'Sales Amount': f'{prev_year}_MTD_Sales'})
    ytd_current_agg = ytd_current_agg.rename(columns={'Sales Amount': f'{target_year}_YTD_Sales'})
    ytd_prev_agg = ytd_prev_agg.rename(columns={'Sales Amount': f'{prev_year}_YTD_Sales'})
    
    # Start with MTD current
    customer_metrics = mtd_current_agg.copy()
    
    # Merge MTD previous
    customer_metrics = customer_metrics.merge(
        mtd_prev_agg[['Customer #', f'{prev_year}_MTD_Sales']],
        on='Customer #',
        how='outer'
    )
    
    # Merge YTD current
    customer_metrics = customer_metrics.merge(
        ytd_current_agg[['Customer #', f'{target_year}_YTD_Sales']],
        on='Customer #',
        how='outer'
    )
    
    # Merge YTD previous
    customer_metrics = customer_metrics.merge(
        ytd_prev_agg[['Customer #', f'{prev_year}_YTD_Sales']],
        on='Customer #',
        how='outer'
    )
    
    # Fill NaN values with 0 for sales columns
    sales_cols = [f'{target_year}_MTD_Sales', f'{prev_year}_MTD_Sales', 
                  f'{target_year}_YTD_Sales', f'{prev_year}_YTD_Sales']
    for col in sales_cols:
        if col in customer_metrics.columns:
            customer_metrics[col] = customer_metrics[col].fillna(0)
    
    # Get complete customer info
    all_customers = pd.concat([
        df_sales[['Customer #', 'Customer Name', channel_col] + 
                 ([sales_rep_col] if sales_rep_col else [])]
    ]).drop_duplicates(subset=['Customer #'])
    
    # Update missing customer info
    cols_to_drop = ['Customer Name', channel_col]
    if sales_rep_col and sales_rep_col in customer_metrics.columns:
        cols_to_drop.append(sales_rep_col)
    customer_metrics = customer_metrics.drop(columns=cols_to_drop, errors='ignore')
    customer_metrics = customer_metrics.merge(all_customers, on='Customer #', how='left')
    
    # Calculate % Achieved
    customer_metrics['MTD_Achieved_%'] = customer_metrics.apply(
        lambda x: (x[f'{target_year}_MTD_Sales'] / x[f'{prev_year}_MTD_Sales'] * 100) 
        if x[f'{prev_year}_MTD_Sales'] > 0 else 0, axis=1
    )
    customer_metrics['YTD_Achieved_%'] = customer_metrics.apply(
        lambda x: (x[f'{target_year}_YTD_Sales'] / x[f'{prev_year}_YTD_Sales'] * 100) 
        if x[f'{prev_year}_YTD_Sales'] > 0 else 0, axis=1
    )
    
    # ============================================================
    # STEP 4: Calculate Channel-Level MTD/YTD Totals
    # ============================================================
    channel_metrics = customer_metrics.groupby(channel_col).agg({
        f'{target_year}_MTD_Sales': 'sum',
        f'{prev_year}_MTD_Sales': 'sum',
        f'{target_year}_YTD_Sales': 'sum',
        f'{prev_year}_YTD_Sales': 'sum',
        'Customer #': 'count'
    }).reset_index()
    
    channel_metrics = channel_metrics.rename(columns={'Customer #': 'Customer_Count'})
    
    # Calculate % Achieved for channels
    channel_metrics['MTD_Achieved_%'] = channel_metrics.apply(
        lambda x: (x[f'{target_year}_MTD_Sales'] / x[f'{prev_year}_MTD_Sales'] * 100) 
        if x[f'{prev_year}_MTD_Sales'] > 0 else 0, axis=1
    )
    channel_metrics['YTD_Achieved_%'] = channel_metrics.apply(
        lambda x: (x[f'{target_year}_YTD_Sales'] / x[f'{prev_year}_YTD_Sales'] * 100) 
        if x[f'{prev_year}_YTD_Sales'] > 0 else 0, axis=1
    )
    
    # Calculate YoY Growth
    channel_metrics['MTD_YoY_Growth'] = channel_metrics.apply(
        lambda x: ((x[f'{target_year}_MTD_Sales'] - x[f'{prev_year}_MTD_Sales']) / x[f'{prev_year}_MTD_Sales'] * 100) 
        if x[f'{prev_year}_MTD_Sales'] > 0 else 0, axis=1
    )
    channel_metrics['YTD_YoY_Growth'] = channel_metrics.apply(
        lambda x: ((x[f'{target_year}_YTD_Sales'] - x[f'{prev_year}_YTD_Sales']) / x[f'{prev_year}_YTD_Sales'] * 100) 
        if x[f'{prev_year}_YTD_Sales'] > 0 else 0, axis=1
    )
    
    # Calculate Grand Total
    grand_total = {
        'Channel': 'GRAND TOTAL',
        'Customer_Count': channel_metrics['Customer_Count'].sum(),
        f'{target_year}_MTD_Sales': channel_metrics[f'{target_year}_MTD_Sales'].sum(),
        f'{prev_year}_MTD_Sales': channel_metrics[f'{prev_year}_MTD_Sales'].sum(),
        f'{target_year}_YTD_Sales': channel_metrics[f'{target_year}_YTD_Sales'].sum(),
        f'{prev_year}_YTD_Sales': channel_metrics[f'{prev_year}_YTD_Sales'].sum(),
    }
    grand_total['MTD_Achieved_%'] = (grand_total[f'{target_year}_MTD_Sales'] / grand_total[f'{prev_year}_MTD_Sales'] * 100) if grand_total[f'{prev_year}_MTD_Sales'] > 0 else 0
    grand_total['YTD_Achieved_%'] = (grand_total[f'{target_year}_YTD_Sales'] / grand_total[f'{prev_year}_YTD_Sales'] * 100) if grand_total[f'{prev_year}_YTD_Sales'] > 0 else 0
    
    # ============================================================
    # STEP 5: Generate Top 10 by Channel for MTD and YTD
    # ============================================================
    channels = customer_metrics[channel_col].dropna().unique()
    channels = [c for c in channels if c != 'Unknown']
    channels = sorted(channels)
    
    results = {
        'channel_metrics': channel_metrics,
        'grand_total': grand_total,
        'top10_by_channel': {},
        'channel_col': channel_col,
        'sales_rep_col': sales_rep_col,
        'target_year': target_year,
        'prev_year': prev_year,
        'channels': channels
    }
    
    for channel in channels:
        channel_data = customer_metrics[customer_metrics[channel_col] == channel].copy()
        
        if channel_data.empty:
            continue
        
        # Top 10 MTD
        top10_mtd = channel_data.nlargest(10, f'{target_year}_MTD_Sales')
        
        # Top 10 YTD
        top10_ytd = channel_data.nlargest(10, f'{target_year}_YTD_Sales')
        
        results['top10_by_channel'][channel] = {
            'MTD': top10_mtd,
            'YTD': top10_ytd
        }
    
    return results


def create_channel_customer_excel_report(results, target_month, target_year, month_name):
    """Create Excel report with channel summary and top 10 customers per channel"""
    
    if results is None:
        return None
    
    prev_year = target_year - 1
    channel_col = results['channel_col']
    sales_rep_col = results['sales_rep_col']
    channel_metrics = results['channel_metrics']
    grand_total = results['grand_total']
    channels = results['channels']
    
    wb = Workbook()
    wb.remove(wb.active)
    
    # Styling
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=10)
    title_font = Font(bold=True, size=14)
    subtitle_font = Font(bold=True, size=12)
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # Color fills
    current_year_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    prev_year_fill = PatternFill(start_color="FCE4D6", end_color="FCE4D6", fill_type="solid")
    achieved_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
    total_fill = PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
    
    # ============================================================
    # Create Summary Sheet with Channel MTD/YTD Totals
    # ============================================================
    ws_summary = wb.create_sheet(title="Channel Summary", index=0)
    
    # Title
    ws_summary.merge_cells('A1:J1')
    ws_summary['A1'] = f"Channel Performance Summary - MTD/YTD ({target_month}/{target_year})"
    ws_summary['A1'].font = Font(bold=True, size=16)
    ws_summary['A1'].alignment = Alignment(horizontal='center')
    
    # Channel Summary Headers
    summary_headers = [
        'Channel', 'Customers',
        f'{target_year} MTD', f'{prev_year} MTD', 'MTD %', 'MTD YoY Growth',
        f'{target_year} YTD', f'{prev_year} YTD', 'YTD %', 'YTD YoY Growth'
    ]
    
    for col_idx, header in enumerate(summary_headers, start=1):
        cell = ws_summary.cell(row=3, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = border
    
    # Write Channel Data
    row_idx = 4
    for _, ch_row in channel_metrics.iterrows():
        if ch_row[channel_col] == 'Unknown':
            continue
        
        ws_summary.cell(row=row_idx, column=1, value=ch_row[channel_col]).border = border
        ws_summary.cell(row=row_idx, column=2, value=ch_row['Customer_Count']).border = border
        
        # MTD Current Year
        cell = ws_summary.cell(row=row_idx, column=3, value=ch_row[f'{target_year}_MTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = current_year_fill
        cell.border = border
        
        # MTD Previous Year
        cell = ws_summary.cell(row=row_idx, column=4, value=ch_row[f'{prev_year}_MTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = prev_year_fill
        cell.border = border
        
        # MTD Achieved %
        cell = ws_summary.cell(row=row_idx, column=5, value=ch_row['MTD_Achieved_%'] / 100)
        cell.number_format = '0.0%'
        cell.fill = achieved_fill
        cell.border = border
        
        # MTD YoY Growth
        cell = ws_summary.cell(row=row_idx, column=6, value=ch_row['MTD_YoY_Growth'] / 100)
        cell.number_format = '+0.0%;-0.0%'
        cell.border = border
        
        # YTD Current Year
        cell = ws_summary.cell(row=row_idx, column=7, value=ch_row[f'{target_year}_YTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = current_year_fill
        cell.border = border
        
        # YTD Previous Year
        cell = ws_summary.cell(row=row_idx, column=8, value=ch_row[f'{prev_year}_YTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = prev_year_fill
        cell.border = border
        
        # YTD Achieved %
        cell = ws_summary.cell(row=row_idx, column=9, value=ch_row['YTD_Achieved_%'] / 100)
        cell.number_format = '0.0%'
        cell.fill = achieved_fill
        cell.border = border
        
        # YTD YoY Growth
        cell = ws_summary.cell(row=row_idx, column=10, value=ch_row['YTD_YoY_Growth'] / 100)
        cell.number_format = '+0.0%;-0.0%'
        cell.border = border
        
        row_idx += 1
    
    # Grand Total Row
    ws_summary.cell(row=row_idx, column=1, value='GRAND TOTAL').font = Font(bold=True)
    ws_summary.cell(row=row_idx, column=1).fill = total_fill
    ws_summary.cell(row=row_idx, column=1).border = border
    
    ws_summary.cell(row=row_idx, column=2, value=grand_total['Customer_Count']).font = Font(bold=True)
    ws_summary.cell(row=row_idx, column=2).fill = total_fill
    ws_summary.cell(row=row_idx, column=2).border = border
    
    cell = ws_summary.cell(row=row_idx, column=3, value=grand_total[f'{target_year}_MTD_Sales'])
    cell.number_format = '$#,##0.00'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    cell = ws_summary.cell(row=row_idx, column=4, value=grand_total[f'{prev_year}_MTD_Sales'])
    cell.number_format = '$#,##0.00'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    cell = ws_summary.cell(row=row_idx, column=5, value=grand_total['MTD_Achieved_%'] / 100)
    cell.number_format = '0.0%'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    mtd_growth = ((grand_total[f'{target_year}_MTD_Sales'] - grand_total[f'{prev_year}_MTD_Sales']) / grand_total[f'{prev_year}_MTD_Sales'] * 100) if grand_total[f'{prev_year}_MTD_Sales'] > 0 else 0
    cell = ws_summary.cell(row=row_idx, column=6, value=mtd_growth / 100)
    cell.number_format = '+0.0%;-0.0%'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    cell = ws_summary.cell(row=row_idx, column=7, value=grand_total[f'{target_year}_YTD_Sales'])
    cell.number_format = '$#,##0.00'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    cell = ws_summary.cell(row=row_idx, column=8, value=grand_total[f'{prev_year}_YTD_Sales'])
    cell.number_format = '$#,##0.00'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    cell = ws_summary.cell(row=row_idx, column=9, value=grand_total['YTD_Achieved_%'] / 100)
    cell.number_format = '0.0%'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    ytd_growth = ((grand_total[f'{target_year}_YTD_Sales'] - grand_total[f'{prev_year}_YTD_Sales']) / grand_total[f'{prev_year}_YTD_Sales'] * 100) if grand_total[f'{prev_year}_YTD_Sales'] > 0 else 0
    cell = ws_summary.cell(row=row_idx, column=10, value=ytd_growth / 100)
    cell.number_format = '+0.0%;-0.0%'
    cell.font = Font(bold=True)
    cell.fill = total_fill
    cell.border = border
    
    # Adjust column widths
    ws_summary.column_dimensions['A'].width = 18
    ws_summary.column_dimensions['B'].width = 12
    ws_summary.column_dimensions['C'].width = 16
    ws_summary.column_dimensions['D'].width = 16
    ws_summary.column_dimensions['E'].width = 10
    ws_summary.column_dimensions['F'].width = 14
    ws_summary.column_dimensions['G'].width = 16
    ws_summary.column_dimensions['H'].width = 16
    ws_summary.column_dimensions['I'].width = 10
    ws_summary.column_dimensions['J'].width = 14
    
    ws_summary.row_dimensions[3].height = 30
    
    # ============================================================
    # Add Charts to Summary Sheet
    # ============================================================
    chart_start_row = row_idx + 3
    
    # Create charts
    fig, axes = plt.subplots(2, 2, figsize=(14, 10))
    fig.suptitle(f'Channel Performance Analysis - {target_month}/{target_year}', fontsize=14, fontweight='bold')
    
    # Filter out Unknown channel for charts
    chart_data = channel_metrics[channel_metrics[channel_col] != 'Unknown'].copy()
    chart_channels = chart_data[channel_col].tolist()
    
    if len(chart_channels) > 0:
        x = np.arange(len(chart_channels))
        width = 0.35
        
        # Chart 1: MTD Sales by Channel
        axes[0, 0].bar(x - width/2, chart_data[f'{target_year}_MTD_Sales'], width, 
                               label=str(target_year), color='#2E86AB', alpha=0.8)
        axes[0, 0].bar(x + width/2, chart_data[f'{prev_year}_MTD_Sales'], width, 
                               label=str(prev_year), color='#A23B72', alpha=0.8)
        axes[0, 0].set_title('MTD Sales by Channel', fontsize=11, fontweight='bold')
        axes[0, 0].set_ylabel('Sales Amount ($)')
        axes[0, 0].set_xticks(x)
        axes[0, 0].set_xticklabels(chart_channels, rotation=45, ha='right', fontsize=9)
        axes[0, 0].legend(fontsize=9)
        axes[0, 0].yaxis.set_major_formatter(plt.FuncFormatter(lambda v, p: f'${v/1000:.0f}K'))
        axes[0, 0].grid(axis='y', alpha=0.3)
        
        # Chart 2: YTD Sales by Channel
        axes[0, 1].bar(x - width/2, chart_data[f'{target_year}_YTD_Sales'], width, 
                               label=str(target_year), color='#2E86AB', alpha=0.8)
        axes[0, 1].bar(x + width/2, chart_data[f'{prev_year}_YTD_Sales'], width, 
                               label=str(prev_year), color='#A23B72', alpha=0.8)
        axes[0, 1].set_title('YTD Sales by Channel', fontsize=11, fontweight='bold')
        axes[0, 1].set_ylabel('Sales Amount ($)')
        axes[0, 1].set_xticks(x)
        axes[0, 1].set_xticklabels(chart_channels, rotation=45, ha='right', fontsize=9)
        axes[0, 1].legend(fontsize=9)
        axes[0, 1].yaxis.set_major_formatter(plt.FuncFormatter(lambda v, p: f'${v/1000:.0f}K'))
        axes[0, 1].grid(axis='y', alpha=0.3)
        
        # Chart 3: MTD % Achieved by Channel
        colors = ['#28a745' if v >= 100 else '#dc3545' for v in chart_data['MTD_Achieved_%']]
        axes[1, 0].barh(chart_channels, chart_data['MTD_Achieved_%'], color=colors, alpha=0.8)
        axes[1, 0].axvline(x=100, color='black', linestyle='--', linewidth=1, label='100% Target')
        axes[1, 0].set_title('MTD % Achieved by Channel', fontsize=11, fontweight='bold')
        axes[1, 0].set_xlabel('% Achieved')
        for i, v in enumerate(chart_data['MTD_Achieved_%']):
            axes[1, 0].text(v + 2, i, f'{v:.1f}%', va='center', fontsize=9)
        axes[1, 0].grid(axis='x', alpha=0.3)
        
        # Chart 4: YTD % Achieved by Channel
        colors = ['#28a745' if v >= 100 else '#dc3545' for v in chart_data['YTD_Achieved_%']]
        axes[1, 1].barh(chart_channels, chart_data['YTD_Achieved_%'], color=colors, alpha=0.8)
        axes[1, 1].axvline(x=100, color='black', linestyle='--', linewidth=1, label='100% Target')
        axes[1, 1].set_title('YTD % Achieved by Channel', fontsize=11, fontweight='bold')
        axes[1, 1].set_xlabel('% Achieved')
        for i, v in enumerate(chart_data['YTD_Achieved_%']):
            axes[1, 1].text(v + 2, i, f'{v:.1f}%', va='center', fontsize=9)
        axes[1, 1].grid(axis='x', alpha=0.3)
    
    plt.tight_layout()
    
    # Save chart to buffer
    img_buffer = io.BytesIO()
    fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
    img_buffer.seek(0)
    plt.close(fig)
    
    img = XLImage(img_buffer)
    ws_summary.add_image(img, f'A{chart_start_row}')
    
    # ============================================================
    # Create Individual Channel Sheets with Top 10 Customers
    # ============================================================
    for channel, data in results['top10_by_channel'].items():
        clean_channel = "".join(c if c.isalnum() or c in (' ', '-', '_') else '_' for c in channel)[:30]
        ws = wb.create_sheet(title=clean_channel)
        
        # Get channel totals
        ch_metrics = channel_metrics[channel_metrics[channel_col] == channel].iloc[0]
        
        # ---- Channel Summary Section ----
        ws.merge_cells('A1:G1')
        ws['A1'] = f"Channel: {channel} - Performance Summary ({target_month}/{target_year})"
        ws['A1'].font = title_font
        ws['A1'].alignment = Alignment(horizontal='center')
        
        # Channel Totals Table
        ws['A3'] = 'Metric'
        ws['B3'] = f'{target_year}'
        ws['C3'] = f'{prev_year}'
        ws['D3'] = '% Achieved'
        ws['E3'] = 'YoY Growth'
        
        for col in range(1, 6):
            ws.cell(row=3, column=col).font = header_font
            ws.cell(row=3, column=col).fill = header_fill
            ws.cell(row=3, column=col).border = border
            ws.cell(row=3, column=col).alignment = Alignment(horizontal='center')
        
        # MTD Row
        ws['A4'] = 'MTD Sales'
        ws['A4'].border = border
        
        cell = ws.cell(row=4, column=2, value=ch_metrics[f'{target_year}_MTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = current_year_fill
        cell.border = border
        
        cell = ws.cell(row=4, column=3, value=ch_metrics[f'{prev_year}_MTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = prev_year_fill
        cell.border = border
        
        cell = ws.cell(row=4, column=4, value=ch_metrics['MTD_Achieved_%'] / 100)
        cell.number_format = '0.0%'
        cell.fill = achieved_fill
        cell.border = border
        
        cell = ws.cell(row=4, column=5, value=ch_metrics['MTD_YoY_Growth'] / 100)
        cell.number_format = '+0.0%;-0.0%'
        cell.border = border
        
        # YTD Row
        ws['A5'] = 'YTD Sales'
        ws['A5'].border = border
        
        cell = ws.cell(row=5, column=2, value=ch_metrics[f'{target_year}_YTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = current_year_fill
        cell.border = border
        
        cell = ws.cell(row=5, column=3, value=ch_metrics[f'{prev_year}_YTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = prev_year_fill
        cell.border = border
        
        cell = ws.cell(row=5, column=4, value=ch_metrics['YTD_Achieved_%'] / 100)
        cell.number_format = '0.0%'
        cell.fill = achieved_fill
        cell.border = border
        
        cell = ws.cell(row=5, column=5, value=ch_metrics['YTD_YoY_Growth'] / 100)
        cell.number_format = '+0.0%;-0.0%'
        cell.border = border
        
        # ---- MTD Top 10 Section ----
        ws.merge_cells('A8:G8')
        ws['A8'] = "Top 10 MTD Customers"
        ws['A8'].font = subtitle_font
        ws['A8'].alignment = Alignment(horizontal='center')
        
        mtd_headers = ['Rank', 'Customer #', 'Customer Name', 
                       f'{target_year} MTD Sales', f'{prev_year} MTD Sales', '% Achieved']
        if sales_rep_col:
            mtd_headers.insert(3, 'Sales Rep')
        
        for col_idx, header in enumerate(mtd_headers, start=1):
            cell = ws.cell(row=9, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            cell.border = border
        
        mtd_df = data['MTD']
        for row_idx_inner, (_, row) in enumerate(mtd_df.iterrows(), start=10):
            rank = row_idx_inner - 9
            col = 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=rank)
            cell.border = border
            cell.alignment = Alignment(horizontal='center')
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row['Customer #'])
            cell.border = border
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row['Customer Name'])
            cell.border = border
            col += 1
            
            if sales_rep_col:
                cell = ws.cell(row=row_idx_inner, column=col, value=row.get(sales_rep_col, ''))
                cell.border = border
                col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row[f'{target_year}_MTD_Sales'])
            cell.number_format = '$#,##0.00'
            cell.fill = current_year_fill
            cell.border = border
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row[f'{prev_year}_MTD_Sales'])
            cell.number_format = '$#,##0.00'
            cell.fill = prev_year_fill
            cell.border = border
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row['MTD_Achieved_%'] / 100)
            cell.number_format = '0.00%'
            cell.fill = achieved_fill
            cell.border = border
        
        # ---- YTD Top 10 Section ----
        ytd_start_row = 10 + len(mtd_df) + 2
        
        ws.merge_cells(f'A{ytd_start_row}:G{ytd_start_row}')
        ws.cell(row=ytd_start_row, column=1, value="Top 10 YTD Customers")
        ws.cell(row=ytd_start_row, column=1).font = subtitle_font
        ws.cell(row=ytd_start_row, column=1).alignment = Alignment(horizontal='center')
        
        ytd_headers = ['Rank', 'Customer #', 'Customer Name', 
                       f'{target_year} YTD Sales', f'{prev_year} YTD Sales', '% Achieved']
        if sales_rep_col:
            ytd_headers.insert(3, 'Sales Rep')
        
        header_row = ytd_start_row + 1
        for col_idx, header in enumerate(ytd_headers, start=1):
            cell = ws.cell(row=header_row, column=col_idx, value=header)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            cell.border = border
        
        ytd_df = data['YTD']
        for row_idx_inner, (_, row) in enumerate(ytd_df.iterrows(), start=header_row + 1):
            rank = row_idx_inner - header_row
            col = 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=rank)
            cell.border = border
            cell.alignment = Alignment(horizontal='center')
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row['Customer #'])
            cell.border = border
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row['Customer Name'])
            cell.border = border
            col += 1
            
            if sales_rep_col:
                cell = ws.cell(row=row_idx_inner, column=col, value=row.get(sales_rep_col, ''))
                cell.border = border
                col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row[f'{target_year}_YTD_Sales'])
            cell.number_format = '$#,##0.00'
            cell.fill = current_year_fill
            cell.border = border
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row[f'{prev_year}_YTD_Sales'])
            cell.number_format = '$#,##0.00'
            cell.fill = prev_year_fill
            cell.border = border
            col += 1
            
            cell = ws.cell(row=row_idx_inner, column=col, value=row['YTD_Achieved_%'] / 100)
            cell.number_format = '0.00%'
            cell.fill = achieved_fill
            cell.border = border
        
        # Adjust column widths
        ws.column_dimensions['A'].width = 8
        ws.column_dimensions['B'].width = 14
        ws.column_dimensions['C'].width = 30
        ws.column_dimensions['D'].width = 16
        ws.column_dimensions['E'].width = 16
        ws.column_dimensions['F'].width = 12
        ws.column_dimensions['G'].width = 12
    
    # Save workbook to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    return output


def generate_sales_rep_performance(sales_df_raw, customer_df, target_month, target_year):
    """
    Generate MTD/YTD performance reports for each Sales Rep.
    
    Args:
        sales_df_raw: Raw sales details DataFrame (already read from file)
        customer_df: Customer list DataFrame with Sales Rep column
        target_month: Month for MTD calculation
        target_year: Year for calculations
    
    Returns:
        Dictionary with sales rep performance data including:
        - sales_rep_metrics: DataFrame with MTD/YTD totals by sales rep
        - grand_total: Overall totals
        - sales_rep_col: Name of the sales rep column
    """
    prev_year = target_year - 1
    
    # ============================================================
    # STEP 1: Parse sales data
    # ============================================================
    col_period = 1
    col_sales_amount = 17
    
    data_rows = []
    current_customer = None
    current_customer_name = None
    
    for idx in range(11, len(sales_df_raw)):
        row = sales_df_raw.iloc[idx]
        first_cell = row.iloc[0]
        
        if pd.isna(first_cell):
            continue
        
        first_cell_str = str(first_cell).strip()
        
        if first_cell_str and not first_cell_str.isdigit() and not first_cell_str.startswith('Item'):
            if any(c.isalpha() for c in first_cell_str) and any(c.isdigit() for c in first_cell_str):
                current_customer = first_cell_str
                cust_name = row.iloc[7] if pd.notna(row.iloc[7]) else ''
                current_customer_name = str(cust_name).strip() if cust_name else current_customer
                continue
        
        if first_cell_str.isdigit() and len(first_cell_str) == 4:
            year = int(first_cell_str)
            
            period_cell = row.iloc[col_period]
            month = None
            if pd.notna(period_cell):
                try:
                    month = int(float(period_cell))
                except (ValueError, TypeError):
                    month = None
            
            sales_amount = row.iloc[col_sales_amount]
            
            if pd.notna(sales_amount) and current_customer and month is not None:
                try:
                    sales_value = float(sales_amount)
                    data_rows.append({
                        'Customer #': current_customer,
                        'Customer Name': current_customer_name,
                        'Year': year,
                        'Month': month,
                        'Sales Amount': sales_value
                    })
                except (ValueError, TypeError):
                    pass
    
    df_sales = pd.DataFrame(data_rows)
    
    if df_sales.empty:
        return None
    
    # ============================================================
    # STEP 2: Join with customer list to get Sales Rep
    # ============================================================
    sales_rep_col = None
    for col in customer_df.columns:
        if 'sales' in col.lower() and 'rep' in col.lower():
            sales_rep_col = col
            break
    
    if not sales_rep_col:
        return None
    
    # Create lookup for customer -> sales rep
    customer_lookup = customer_df[['Customer #', sales_rep_col]].drop_duplicates(subset=['Customer #'])
    
    # Merge to add sales rep to sales data
    df_sales = df_sales.merge(customer_lookup, on='Customer #', how='left')
    
    # Fill missing sales reps with 'Unassigned'
    df_sales[sales_rep_col] = df_sales[sales_rep_col].fillna('Unassigned')
    
    # ============================================================
    # STEP 3: Calculate MTD and YTD metrics by Sales Rep
    # ============================================================
    df_current = df_sales[df_sales['Year'] == target_year].copy()
    df_prev = df_sales[df_sales['Year'] == prev_year].copy()
    
    # Calculate MTD (specific month only)
    df_mtd_current = df_current[df_current['Month'] == target_month]
    df_mtd_prev = df_prev[df_prev['Month'] == target_month]
    
    # Calculate YTD (months 1 to target_month)
    df_ytd_current = df_current[df_current['Month'] <= target_month]
    df_ytd_prev = df_prev[df_prev['Month'] <= target_month]
    
    # Aggregate by Sales Rep
    def aggregate_by_rep(df, rep_col):
        return df.groupby(rep_col).agg({
            'Sales Amount': 'sum',
            'Customer #': 'nunique'
        }).reset_index()
    
    mtd_current_agg = aggregate_by_rep(df_mtd_current, sales_rep_col)
    mtd_prev_agg = aggregate_by_rep(df_mtd_prev, sales_rep_col)
    ytd_current_agg = aggregate_by_rep(df_ytd_current, sales_rep_col)
    ytd_prev_agg = aggregate_by_rep(df_ytd_prev, sales_rep_col)
    
    # Rename columns
    mtd_current_agg = mtd_current_agg.rename(columns={
        'Sales Amount': f'{target_year}_MTD_Sales',
        'Customer #': 'Customer_Count'
    })
    mtd_prev_agg = mtd_prev_agg.rename(columns={
        'Sales Amount': f'{prev_year}_MTD_Sales',
        'Customer #': 'Customer_Count_Prev'
    })
    ytd_current_agg = ytd_current_agg.rename(columns={
        'Sales Amount': f'{target_year}_YTD_Sales',
        'Customer #': 'Customer_Count_YTD'
    })
    ytd_prev_agg = ytd_prev_agg.rename(columns={
        'Sales Amount': f'{prev_year}_YTD_Sales',
        'Customer #': 'Customer_Count_YTD_Prev'
    })
    
    # Merge all metrics
    sales_rep_metrics = mtd_current_agg[[sales_rep_col, f'{target_year}_MTD_Sales', 'Customer_Count']].copy()
    
    sales_rep_metrics = sales_rep_metrics.merge(
        mtd_prev_agg[[sales_rep_col, f'{prev_year}_MTD_Sales']],
        on=sales_rep_col, how='outer'
    )
    sales_rep_metrics = sales_rep_metrics.merge(
        ytd_current_agg[[sales_rep_col, f'{target_year}_YTD_Sales']],
        on=sales_rep_col, how='outer'
    )
    sales_rep_metrics = sales_rep_metrics.merge(
        ytd_prev_agg[[sales_rep_col, f'{prev_year}_YTD_Sales']],
        on=sales_rep_col, how='outer'
    )
    
    # Fill NaN with 0
    for col in [f'{target_year}_MTD_Sales', f'{prev_year}_MTD_Sales', 
                f'{target_year}_YTD_Sales', f'{prev_year}_YTD_Sales', 'Customer_Count']:
        if col in sales_rep_metrics.columns:
            sales_rep_metrics[col] = sales_rep_metrics[col].fillna(0)
    
    # Calculate % Achieved
    sales_rep_metrics['MTD_Achieved_%'] = sales_rep_metrics.apply(
        lambda x: (x[f'{target_year}_MTD_Sales'] / x[f'{prev_year}_MTD_Sales'] * 100) 
        if x[f'{prev_year}_MTD_Sales'] > 0 else 0, axis=1
    )
    sales_rep_metrics['YTD_Achieved_%'] = sales_rep_metrics.apply(
        lambda x: (x[f'{target_year}_YTD_Sales'] / x[f'{prev_year}_YTD_Sales'] * 100) 
        if x[f'{prev_year}_YTD_Sales'] > 0 else 0, axis=1
    )
    
    # Calculate YoY Growth
    sales_rep_metrics['MTD_YoY_Growth'] = sales_rep_metrics.apply(
        lambda x: ((x[f'{target_year}_MTD_Sales'] - x[f'{prev_year}_MTD_Sales']) / x[f'{prev_year}_MTD_Sales'] * 100) 
        if x[f'{prev_year}_MTD_Sales'] > 0 else 0, axis=1
    )
    sales_rep_metrics['YTD_YoY_Growth'] = sales_rep_metrics.apply(
        lambda x: ((x[f'{target_year}_YTD_Sales'] - x[f'{prev_year}_YTD_Sales']) / x[f'{prev_year}_YTD_Sales'] * 100) 
        if x[f'{prev_year}_YTD_Sales'] > 0 else 0, axis=1
    )
    
    # Sort by current year YTD sales descending
    sales_rep_metrics = sales_rep_metrics.sort_values(f'{target_year}_YTD_Sales', ascending=False)
    
    # Calculate Grand Total
    grand_total = {
        'Sales Rep': 'GRAND TOTAL',
        'Customer_Count': int(sales_rep_metrics['Customer_Count'].sum()),
        f'{target_year}_MTD_Sales': sales_rep_metrics[f'{target_year}_MTD_Sales'].sum(),
        f'{prev_year}_MTD_Sales': sales_rep_metrics[f'{prev_year}_MTD_Sales'].sum(),
        f'{target_year}_YTD_Sales': sales_rep_metrics[f'{target_year}_YTD_Sales'].sum(),
        f'{prev_year}_YTD_Sales': sales_rep_metrics[f'{prev_year}_YTD_Sales'].sum(),
    }
    grand_total['MTD_Achieved_%'] = (grand_total[f'{target_year}_MTD_Sales'] / grand_total[f'{prev_year}_MTD_Sales'] * 100) if grand_total[f'{prev_year}_MTD_Sales'] > 0 else 0
    grand_total['YTD_Achieved_%'] = (grand_total[f'{target_year}_YTD_Sales'] / grand_total[f'{prev_year}_YTD_Sales'] * 100) if grand_total[f'{prev_year}_YTD_Sales'] > 0 else 0
    
    return {
        'sales_rep_metrics': sales_rep_metrics,
        'grand_total': grand_total,
        'sales_rep_col': sales_rep_col
    }


def create_sales_rep_excel_report(results, target_month, target_year, month_name):
    """
    Create Excel report for Sales Rep performance with charts.
    """
    if results is None:
        return None
    
    prev_year = target_year - 1
    sales_rep_metrics = results['sales_rep_metrics']
    grand_total = results['grand_total']
    sales_rep_col = results['sales_rep_col']
    
    wb = Workbook()
    ws = wb.active
    ws.title = "Sales Rep Summary"
    
    # Styling
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=10)
    border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # Color fills
    current_year_fill = PatternFill(start_color="E2EFDA", end_color="E2EFDA", fill_type="solid")
    prev_year_fill = PatternFill(start_color="FCE4D6", end_color="FCE4D6", fill_type="solid")
    achieved_fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
    total_fill = PatternFill(start_color="DDEBF7", end_color="DDEBF7", fill_type="solid")
    
    # Title
    ws.merge_cells('A1:J1')
    ws['A1'] = f"Sales Rep Performance Summary - MTD/YTD ({month_name} {target_year})"
    ws['A1'].font = Font(bold=True, size=16)
    ws['A1'].alignment = Alignment(horizontal='center')
    
    # Headers
    headers = [
        'Sales Rep', 'Customers',
        f'{target_year} MTD', f'{prev_year} MTD', 'MTD %', 'MTD YoY Growth',
        f'{target_year} YTD', f'{prev_year} YTD', 'YTD %', 'YTD YoY Growth'
    ]
    
    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=3, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
        cell.border = border
    
    # Write Sales Rep Data
    row_idx = 4
    for _, rep_row in sales_rep_metrics.iterrows():
        if rep_row[sales_rep_col] == 'Unassigned':
            continue
        
        ws.cell(row=row_idx, column=1, value=rep_row[sales_rep_col]).border = border
        ws.cell(row=row_idx, column=2, value=int(rep_row['Customer_Count'])).border = border
        
        # MTD Current Year
        cell = ws.cell(row=row_idx, column=3, value=rep_row[f'{target_year}_MTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = current_year_fill
        cell.border = border
        
        # MTD Previous Year
        cell = ws.cell(row=row_idx, column=4, value=rep_row[f'{prev_year}_MTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = prev_year_fill
        cell.border = border
        
        # MTD Achieved %
        cell = ws.cell(row=row_idx, column=5, value=rep_row['MTD_Achieved_%'] / 100)
        cell.number_format = '0.0%'
        cell.fill = achieved_fill
        cell.border = border
        
        # MTD YoY Growth
        cell = ws.cell(row=row_idx, column=6, value=rep_row['MTD_YoY_Growth'] / 100)
        cell.number_format = '+0.0%;-0.0%'
        cell.border = border
        
        # YTD Current Year
        cell = ws.cell(row=row_idx, column=7, value=rep_row[f'{target_year}_YTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = current_year_fill
        cell.border = border
        
        # YTD Previous Year
        cell = ws.cell(row=row_idx, column=8, value=rep_row[f'{prev_year}_YTD_Sales'])
        cell.number_format = '$#,##0.00'
        cell.fill = prev_year_fill
        cell.border = border
        
        # YTD Achieved %
        cell = ws.cell(row=row_idx, column=9, value=rep_row['YTD_Achieved_%'] / 100)
        cell.number_format = '0.0%'
        cell.fill = achieved_fill
        cell.border = border
        
        # YTD YoY Growth
        cell = ws.cell(row=row_idx, column=10, value=rep_row['YTD_YoY_Growth'] / 100)
        cell.number_format = '+0.0%;-0.0%'
        cell.border = border
        
        row_idx += 1
    
    # Grand Total Row
    for col in range(1, 11):
        ws.cell(row=row_idx, column=col).fill = total_fill
        ws.cell(row=row_idx, column=col).border = border
        ws.cell(row=row_idx, column=col).font = Font(bold=True)
    
    ws.cell(row=row_idx, column=1, value='GRAND TOTAL')
    ws.cell(row=row_idx, column=2, value=grand_total['Customer_Count'])
    
    cell = ws.cell(row=row_idx, column=3, value=grand_total[f'{target_year}_MTD_Sales'])
    cell.number_format = '$#,##0.00'
    
    cell = ws.cell(row=row_idx, column=4, value=grand_total[f'{prev_year}_MTD_Sales'])
    cell.number_format = '$#,##0.00'
    
    cell = ws.cell(row=row_idx, column=5, value=grand_total['MTD_Achieved_%'] / 100)
    cell.number_format = '0.0%'
    
    mtd_growth = ((grand_total[f'{target_year}_MTD_Sales'] - grand_total[f'{prev_year}_MTD_Sales']) / grand_total[f'{prev_year}_MTD_Sales'] * 100) if grand_total[f'{prev_year}_MTD_Sales'] > 0 else 0
    cell = ws.cell(row=row_idx, column=6, value=mtd_growth / 100)
    cell.number_format = '+0.0%;-0.0%'
    
    cell = ws.cell(row=row_idx, column=7, value=grand_total[f'{target_year}_YTD_Sales'])
    cell.number_format = '$#,##0.00'
    
    cell = ws.cell(row=row_idx, column=8, value=grand_total[f'{prev_year}_YTD_Sales'])
    cell.number_format = '$#,##0.00'
    
    cell = ws.cell(row=row_idx, column=9, value=grand_total['YTD_Achieved_%'] / 100)
    cell.number_format = '0.0%'
    
    ytd_growth = ((grand_total[f'{target_year}_YTD_Sales'] - grand_total[f'{prev_year}_YTD_Sales']) / grand_total[f'{prev_year}_YTD_Sales'] * 100) if grand_total[f'{prev_year}_YTD_Sales'] > 0 else 0
    cell = ws.cell(row=row_idx, column=10, value=ytd_growth / 100)
    cell.number_format = '+0.0%;-0.0%'
    
    # Adjust column widths
    ws.column_dimensions['A'].width = 20
    ws.column_dimensions['B'].width = 12
    ws.column_dimensions['C'].width = 16
    ws.column_dimensions['D'].width = 16
    ws.column_dimensions['E'].width = 10
    ws.column_dimensions['F'].width = 14
    ws.column_dimensions['G'].width = 16
    ws.column_dimensions['H'].width = 16
    ws.column_dimensions['I'].width = 10
    ws.column_dimensions['J'].width = 14
    
    ws.row_dimensions[3].height = 30
    
    # ============================================================
    # Add Charts
    # ============================================================
    chart_start_row = row_idx + 3
    
    # Filter out Unassigned for charts
    chart_data = sales_rep_metrics[sales_rep_metrics[sales_rep_col] != 'Unassigned'].copy()
    chart_reps = chart_data[sales_rep_col].tolist()
    
    if len(chart_reps) > 0:
        fig, axes = plt.subplots(2, 2, figsize=(14, 10))
        fig.suptitle(f'Sales Rep Performance Analysis - {month_name} {target_year}', fontsize=14, fontweight='bold')
        
        x = np.arange(len(chart_reps))
        width = 0.35
        
        # Chart 1: MTD Sales by Sales Rep
        axes[0, 0].bar(x - width/2, chart_data[f'{target_year}_MTD_Sales'], width, 
                       label=str(target_year), color='#2E86AB', alpha=0.8)
        axes[0, 0].bar(x + width/2, chart_data[f'{prev_year}_MTD_Sales'], width, 
                       label=str(prev_year), color='#A23B72', alpha=0.8)
        axes[0, 0].set_title('MTD Sales by Sales Rep', fontsize=11, fontweight='bold')
        axes[0, 0].set_ylabel('Sales Amount ($)')
        axes[0, 0].set_xticks(x)
        axes[0, 0].set_xticklabels(chart_reps, rotation=45, ha='right', fontsize=9)
        axes[0, 0].legend(fontsize=9)
        axes[0, 0].yaxis.set_major_formatter(plt.FuncFormatter(lambda v, p: f'${v/1000:.0f}K'))
        axes[0, 0].grid(axis='y', alpha=0.3)
        
        # Chart 2: YTD Sales by Sales Rep
        axes[0, 1].bar(x - width/2, chart_data[f'{target_year}_YTD_Sales'], width, 
                       label=str(target_year), color='#2E86AB', alpha=0.8)
        axes[0, 1].bar(x + width/2, chart_data[f'{prev_year}_YTD_Sales'], width, 
                       label=str(prev_year), color='#A23B72', alpha=0.8)
        axes[0, 1].set_title('YTD Sales by Sales Rep', fontsize=11, fontweight='bold')
        axes[0, 1].set_ylabel('Sales Amount ($)')
        axes[0, 1].set_xticks(x)
        axes[0, 1].set_xticklabels(chart_reps, rotation=45, ha='right', fontsize=9)
        axes[0, 1].legend(fontsize=9)
        axes[0, 1].yaxis.set_major_formatter(plt.FuncFormatter(lambda v, p: f'${v/1000:.0f}K'))
        axes[0, 1].grid(axis='y', alpha=0.3)
        
        # Chart 3: MTD % Achieved by Sales Rep
        colors = ['#28a745' if v >= 100 else '#dc3545' for v in chart_data['MTD_Achieved_%']]
        axes[1, 0].barh(chart_reps, chart_data['MTD_Achieved_%'], color=colors, alpha=0.8)
        axes[1, 0].axvline(x=100, color='black', linestyle='--', linewidth=1)
        axes[1, 0].set_title('MTD % Achieved by Sales Rep', fontsize=11, fontweight='bold')
        axes[1, 0].set_xlabel('% Achieved')
        for i, v in enumerate(chart_data['MTD_Achieved_%']):
            axes[1, 0].text(v + 2, i, f'{v:.1f}%', va='center', fontsize=9)
        axes[1, 0].grid(axis='x', alpha=0.3)
        
        # Chart 4: YTD % Achieved by Sales Rep
        colors = ['#28a745' if v >= 100 else '#dc3545' for v in chart_data['YTD_Achieved_%']]
        axes[1, 1].barh(chart_reps, chart_data['YTD_Achieved_%'], color=colors, alpha=0.8)
        axes[1, 1].axvline(x=100, color='black', linestyle='--', linewidth=1)
        axes[1, 1].set_title('YTD % Achieved by Sales Rep', fontsize=11, fontweight='bold')
        axes[1, 1].set_xlabel('% Achieved')
        for i, v in enumerate(chart_data['YTD_Achieved_%']):
            axes[1, 1].text(v + 2, i, f'{v:.1f}%', va='center', fontsize=9)
        axes[1, 1].grid(axis='x', alpha=0.3)
        
        plt.tight_layout()
        
        # Save chart to buffer
        img_buffer = io.BytesIO()
        fig.savefig(img_buffer, format='png', dpi=150, bbox_inches='tight')
        img_buffer.seek(0)
        plt.close(fig)
        
        img = XLImage(img_buffer)
        ws.add_image(img, f'A{chart_start_row}')
    
    # Save workbook to bytes
    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    
    return output


# Main processing
if uploaded_file is not None:
    # Check if all required files are uploaded
    if brand_mapping_file is None:
        st.warning("⚠️ **ItemMaster File Required**: Please upload the ItemMaster Excel file to generate brand reports. This is required because different brands may share the same first 3 letters in item numbers.")
        st.stop()
    
    if sales_details_file is None:
        st.warning("⚠️ **Sales Details File Required**: Please upload the Sales Details Excel file to generate customer channel reports.")
        st.stop()
    
    if customer_list_file is None:
        st.warning("⚠️ **Sales Customer List File Required**: Please upload the Sales Customer List Excel file with Channel column to generate customer channel reports.")
        st.stop()
    
    try:
        with st.spinner("Processing your file..."):
            # Read Excel file
            excel_file = pd.ExcelFile(uploaded_file)
            
            # Success message in a cleaner format
            st.success(f"✅ File uploaded successfully! Found {len(excel_file.sheet_names)} sheet(s)")
            
            # Read and process all sheets
            dataframes = {}
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(excel_file, sheet_name=sheet_name, header=9)
                df = df.dropna(axis=1, how='all')

                unnamed_cols = [col for col in df.columns if isinstance(col, str) and col.startswith('Unnamed')]
                cols_to_drop = [col for col in unnamed_cols if df[col].isna().all()]
                if cols_to_drop:
                    df = df.drop(columns=cols_to_drop)

                if 'Unnamed: 0' in df.columns and 'Item Number' not in df.columns:
                    df = df.rename(columns={'Unnamed: 0': 'Item Number'})
                dataframes[sheet_name] = df
            
            # Split by items
            all_items = {}
            for sheet_name, df in dataframes.items():
                items = split_by_items(df)
                for item_num, item_df in items.items():
                    key = f"{sheet_name}_{item_num}"
                    all_items[key] = item_df
            
            # Info badge
            st.markdown(f"""
            <div style='background-color: #e7f3ff; padding: 0.8rem; border-radius: 8px; margin: 1rem 0;'>
                📦 <b>Extracted {len(all_items)} items</b> from {len(excel_file.sheet_names)} sheet(s)
            </div>
            """, unsafe_allow_html=True)
            
            # Load brand mapping if brand mapping file is uploaded
            brand_mapping = {}
            brand_normalized_map = {}
            if brand_mapping_file is not None:
                brand_mapping, brand_normalized_map = load_brand_mapping(brand_mapping_file)
            
            # Calculate metrics
            results_current = calculate_sales_metrics(all_items, target_month, target_year)
            results_previous = calculate_sales_metrics(all_items, target_month, comparison_year)
            
            # Display results with improved layout
            st.markdown("---")
            st.markdown(f"## � Performance Overview - {month_name} {target_year}")
            
            # Key metrics in cards
            col1, col2, col3, col4 = st.columns(4)
            
            mtd_change = ((results_current['MTD Gross Sales'] / results_previous['MTD Gross Sales'] - 1) * 100) if results_previous['MTD Gross Sales'] > 0 else 0
            ytd_change = ((results_current['YTD Gross Sales'] / results_previous['YTD Gross Sales'] - 1) * 100) if results_previous['YTD Gross Sales'] > 0 else 0
            mtd_gp_change = results_current['MTD GP%'] - results_previous['MTD GP%']
            ytd_gp_change = results_current['YTD GP%'] - results_previous['YTD GP%']
            
            with col1:
                st.metric(
                    "💰 MTD Gross Sales", 
                    f"${results_current['MTD Gross Sales']:,.0f}",
                    f"{mtd_change:+.1f}%" if mtd_change != 0 else "0%",
                    delta_color="normal"
                )
            
            with col2:
                st.metric(
                    "📊 MTD GP%", 
                    f"{results_current['MTD GP%']:.2f}%",
                    f"{mtd_gp_change:+.2f}%",
                    delta_color="normal"
                )
            
            with col3:
                st.metric(
                    "💵 YTD Gross Sales", 
                    f"${results_current['YTD Gross Sales']:,.0f}",
                    f"{ytd_change:+.1f}%" if ytd_change != 0 else "0%",
                    delta_color="normal"
                )
            
            with col4:
                st.metric(
                    "📈 YTD GP%", 
                    f"{results_current['YTD GP%']:.2f}%",
                    f"{ytd_gp_change:+.2f}%",
                    delta_color="normal"
                )
            
            # Detailed comparison table in expander
            with st.expander("📋 View Detailed Comparison Table", expanded=False):
                mtd_achieved = (results_current['MTD Gross Sales'] / results_previous['MTD Gross Sales'] * 100) if results_previous['MTD Gross Sales'] > 0 else 0
                ytd_achieved = (results_current['YTD Gross Sales'] / results_previous['YTD Gross Sales'] * 100) if results_previous['YTD Gross Sales'] > 0 else 0
                mtd_gp_achieved = "0%" if results_previous['MTD GP%'] == 0 else f"{(results_current['MTD GP%'] / results_previous['MTD GP%'] * 100):.0f}%"
                ytd_gp_achieved = "0%" if results_previous['YTD GP%'] == 0 else f"{(results_current['YTD GP%'] / results_previous['YTD GP%'] * 100):.0f}%"
                
                summary_data = {
                    'Period': [str(target_year), str(comparison_year), '% Achieved', f'{target_year} Budget', '% vs Budget'],
                    'MTD Gross Sales': [
                        f"$ {results_current['MTD Gross Sales']:,.2f}",
                        f"$ {results_previous['MTD Gross Sales']:,.2f}",
                        f"{mtd_achieved:.0f}%",
                        "",
                        "0%"
                    ],
                    'MTD GP%': [
                        f"{results_current['MTD GP%']:.2f}%",
                        f"{results_previous['MTD GP%']:.2f}%",
                        mtd_gp_achieved,
                        "",
                        "0%"
                    ],
                    'YTD Gross Sales': [
                        f"$ {results_current['YTD Gross Sales']:,.2f}",
                        f"$ {results_previous['YTD Gross Sales']:,.2f}",
                        f"{ytd_achieved:.0f}%",
                        "",
                        "0%"
                    ],
                    'YTD GP%': [
                        f"{results_current['YTD GP%']:.2f}%",
                        f"{results_previous['YTD GP%']:.2f}%",
                        ytd_gp_achieved,
                        "",
                        "0%"
                    ]
                }
                
                summary_df = pd.DataFrame(summary_data)
                st.dataframe(summary_df, use_container_width=True, hide_index=True)
            
            # Generate all reports first (but don't show download section yet)
            excel_output = create_excel_report(
                results_current, 
                results_previous, 
                target_month, 
                target_year, 
                comparison_year,
                month_name
            )
            
            brand_mtd_output, brand_mtd_df, total_brands = create_brand_report(
                all_items, target_month, target_year, comparison_year, month_name, 'MTD', brand_mapping, brand_normalized_map
            )
            
            brand_ytd_output, brand_ytd_df, _ = create_brand_report(
                all_items, target_month, target_year, comparison_year, month_name, 'YTD', brand_mapping, brand_normalized_map
            )
            
            sku_mtd_output, sku_mtd_df = create_sku_report(
                all_items, target_month, target_year, comparison_year, month_name, 'MTD'
            )
            
            sku_ytd_output, sku_ytd_df = create_sku_report(
                all_items, target_month, target_year, comparison_year, month_name, 'YTD'
            )
            
            # Data preview section
            st.markdown("---")
            st.markdown("## 🔍 Data Preview")
            
            if brand_mtd_output:
                st.markdown("### 🏆 Top 10 Brands")
                
                tab1, tab2 = st.tabs(["📈 MTD Performance", "📊 YTD Performance"])
                
                with tab1:
                    preview_mtd = []
                    for idx, row in brand_mtd_df.reset_index(drop=True).iterrows():
                        preview_mtd.append({
                            'Rank': idx + 1,
                            'Brand': row['brand_name'],
                            f'{target_year} MTD Sales': f"${row[f'{target_year}_mtd_sales']:,.2f}",
                            f'{target_year} GP%': f"{row[f'{target_year}_mtd_gp']:.2f}%",
                            f'{comparison_year} MTD Sales': f"${row[f'{comparison_year}_mtd_sales']:,.2f}",
                            f'{comparison_year} GP%': f"{row[f'{comparison_year}_mtd_gp']:.2f}%",
                            '% Achieved': f"{row['mtd_achieved_pct']:.2f}%"
                        })
                    
                    preview_mtd_df = pd.DataFrame(preview_mtd)
                    st.dataframe(preview_mtd_df, use_container_width=True, hide_index=True)
                
                with tab2:
                    if brand_ytd_output:
                        preview_ytd = []
                        for idx, row in brand_ytd_df.reset_index(drop=True).iterrows():
                            preview_ytd.append({
                                'Rank': idx + 1,
                                'Brand': row['brand_name'],
                                f'{target_year} YTD Sales': f"${row[f'{target_year}_ytd_sales']:,.2f}",
                                f'{target_year} GP%': f"{row[f'{target_year}_ytd_gp']:.2f}%",
                                f'{comparison_year} YTD Sales': f"${row[f'{comparison_year}_ytd_sales']:,.2f}",
                                f'{comparison_year} GP%': f"{row[f'{comparison_year}_ytd_gp']:.2f}%",
                                '% Achieved': f"{row['ytd_achieved_pct']:.2f}%"
                            })
                        
                        preview_ytd_df = pd.DataFrame(preview_ytd)
                        st.dataframe(preview_ytd_df, use_container_width=True, hide_index=True)
            
            # SKU previews
            if sku_mtd_output:
                st.markdown("### 🏷️ Top 20 SKUs")
                
                tab3, tab4 = st.tabs(["📈 MTD SKU Performance", "📊 YTD SKU Performance"])
                
                with tab3:
                    preview_sku_mtd = []
                    for idx, row in sku_mtd_df.reset_index(drop=True).iterrows():
                        preview_sku_mtd.append({
                            'Rank': idx + 1,
                            'Code': row['code'],
                            'Item Name': row['name'][:50] + '...' if len(row['name']) > 50 else row['name'],
                            f'{target_year} MTD Sales': f"${row[f'{target_year}_mtd_sales']:,.2f}",
                            f'{target_year} GP%': f"{row[f'{target_year}_mtd_gp']:.2f}%",
                            f'{comparison_year} MTD Sales': f"${row[f'{comparison_year}_mtd_sales']:,.2f}",
                            f'{comparison_year} GP%': f"{row[f'{comparison_year}_mtd_gp']:.2f}%"
                        })
                    
                    preview_sku_mtd_df = pd.DataFrame(preview_sku_mtd)
                    st.dataframe(preview_sku_mtd_df, use_container_width=True, hide_index=True)
                
                with tab4:
                    if sku_ytd_output:
                        preview_sku_ytd = []
                        for idx, row in sku_ytd_df.reset_index(drop=True).iterrows():
                            preview_sku_ytd.append({
                                'Rank': idx + 1,
                                'Code': row['code'],
                                'Item Name': row['name'][:50] + '...' if len(row['name']) > 50 else row['name'],
                                f'{target_year} YTD Sales': f"${row[f'{target_year}_ytd_sales']:,.2f}",
                                f'{target_year} GP%': f"{row[f'{target_year}_ytd_gp']:.2f}%",
                                f'{comparison_year} YTD Sales': f"${row[f'{comparison_year}_ytd_sales']:,.2f}",
                                f'{comparison_year} GP%': f"{row[f'{comparison_year}_ytd_gp']:.2f}%"
                            })
                        
                        preview_sku_ytd_df = pd.DataFrame(preview_sku_ytd)
                        st.dataframe(preview_sku_ytd_df, use_container_width=True, hide_index=True)
            
            # Customer Channel Performance Preview
            st.markdown("### 👥 Customer Channel Performance")
            
            try:
                channel_results = generate_top10_customers_by_channel(
                    sales_details_file,
                    customer_list_file,
                    target_month,
                    target_year
                )
                
                if channel_results is not None:
                    channel_excel_output = create_channel_customer_excel_report(
                        channel_results, target_month, target_year, month_name
                    )
                    
                    # Display channel summary
                    channel_col = channel_results['channel_col']
                    channel_metrics = channel_results['channel_metrics']
                    prev_year = target_year - 1
                    
                    channel_preview = []
                    for _, ch_row in channel_metrics.iterrows():
                        if ch_row[channel_col] != 'Unknown':
                            channel_preview.append({
                                'Channel': ch_row[channel_col],
                                'Customers': ch_row['Customer_Count'],
                                f'{target_year} MTD': f"${ch_row[f'{target_year}_MTD_Sales']:,.2f}",
                                f'{prev_year} MTD': f"${ch_row[f'{prev_year}_MTD_Sales']:,.2f}",
                                'MTD %': f"{ch_row['MTD_Achieved_%']:.1f}%",
                                f'{target_year} YTD': f"${ch_row[f'{target_year}_YTD_Sales']:,.2f}",
                                f'{prev_year} YTD': f"${ch_row[f'{prev_year}_YTD_Sales']:,.2f}",
                                'YTD %': f"{ch_row['YTD_Achieved_%']:.1f}%"
                            })
                    
                    if channel_preview:
                        channel_df_display = pd.DataFrame(channel_preview)
                        st.dataframe(channel_df_display, use_container_width=True, hide_index=True)
                    
                    # Grand Total
                    gt = channel_results['grand_total']
                    mtd_sales_fmt = f"${gt[f'{target_year}_MTD_Sales']:,.2f}"
                    ytd_sales_fmt = f"${gt[f'{target_year}_YTD_Sales']:,.2f}"
                    st.markdown(f"""
                    **Grand Total:** {gt['Customer_Count']} customers | 
                    MTD: {mtd_sales_fmt} ({gt['MTD_Achieved_%']:.1f}% vs {prev_year}) | 
                    YTD: {ytd_sales_fmt} ({gt['YTD_Achieved_%']:.1f}% vs {prev_year})
                    """)
                    
                    st.success(f"✅ Channel report includes {len(channel_results['channels'])} channels with Top 10 MTD/YTD customers each")
                else:
                    channel_excel_output = None
                    channel_results = None
                    st.error("❌ Could not process channel customer data")
            except Exception as e:
                channel_excel_output = None
                channel_results = None
                st.error(f"❌ Error generating channel report: {str(e)}")
            
            # Sales Rep Performance Preview
            st.markdown("### 👤 Sales Rep Performance")
            
            try:
                # Read raw sales details for sales rep analysis
                sales_details_df_raw = pd.read_excel(sales_details_file, header=None)
                customer_df_for_rep = pd.read_excel(customer_list_file)
                
                sales_rep_results = generate_sales_rep_performance(
                    sales_details_df_raw,
                    customer_df_for_rep,
                    target_month,
                    target_year
                )
                
                if sales_rep_results is not None:
                    sales_rep_excel_output = create_sales_rep_excel_report(
                        sales_rep_results, target_month, target_year, month_name
                    )
                    
                    # Display sales rep summary
                    rep_col = sales_rep_results['sales_rep_col']
                    rep_metrics = sales_rep_results['sales_rep_metrics']
                    prev_year = target_year - 1
                    
                    rep_preview = []
                    for _, rep_row in rep_metrics.iterrows():
                        if rep_row[rep_col] != 'Unassigned':
                            rep_preview.append({
                                'Sales Rep': rep_row[rep_col],
                                'Customers': int(rep_row['Customer_Count']),
                                f'{target_year} MTD': f"${rep_row[f'{target_year}_MTD_Sales']:,.2f}",
                                f'{prev_year} MTD': f"${rep_row[f'{prev_year}_MTD_Sales']:,.2f}",
                                'MTD %': f"{rep_row['MTD_Achieved_%']:.1f}%",
                                f'{target_year} YTD': f"${rep_row[f'{target_year}_YTD_Sales']:,.2f}",
                                f'{prev_year} YTD': f"${rep_row[f'{prev_year}_YTD_Sales']:,.2f}",
                                'YTD %': f"{rep_row['YTD_Achieved_%']:.1f}%"
                            })
                    
                    if rep_preview:
                        rep_df_display = pd.DataFrame(rep_preview)
                        st.dataframe(rep_df_display, use_container_width=True, hide_index=True)
                    
                    # Grand Total
                    gt_rep = sales_rep_results['grand_total']
                    rep_mtd_sales_fmt = f"${gt_rep[f'{target_year}_MTD_Sales']:,.2f}"
                    rep_ytd_sales_fmt = f"${gt_rep[f'{target_year}_YTD_Sales']:,.2f}"
                    st.markdown(f"""
                    **Grand Total:** {gt_rep['Customer_Count']} customers | 
                    MTD: {rep_mtd_sales_fmt} ({gt_rep['MTD_Achieved_%']:.1f}% vs {prev_year}) | 
                    YTD: {rep_ytd_sales_fmt} ({gt_rep['YTD_Achieved_%']:.1f}% vs {prev_year})
                    """)
                    
                    num_reps = len([r for _, r in rep_metrics.iterrows() if r[rep_col] != 'Unassigned'])
                    st.success(f"✅ Sales Rep report includes {num_reps} sales representatives with MTD/YTD analysis")
                else:
                    sales_rep_excel_output = None
                    sales_rep_results = None
                    st.error("❌ Could not process sales rep data. Make sure the customer list has a 'Sales Rep' column.")
            except Exception as e:
                sales_rep_excel_output = None
                sales_rep_results = None
                st.error(f"❌ Error generating sales rep report: {str(e)}")
            
            # Generate PowerPoint presentation (after channel and sales rep data is available)
            ppt_output = create_powerpoint_presentation(
                results_current, results_previous, target_month, target_year,
                comparison_year, month_name, brand_mtd_df, brand_ytd_df,
                sku_mtd_df, sku_ytd_df, 
                channel_results=channel_results,
                sales_rep_results=sales_rep_results
            )
            
            # Additional info
            st.info(f"ℹ️ Processed {results_current['items_processed']} items with data for {target_year} | {total_brands} brands identified")
            
            # Download reports section with better organization
            st.markdown("---")
            st.markdown("## 📥 Download Reports")
            
            # Organized download tabs
            tab1, tab2, tab3, tab4, tab5, tab6 = st.tabs(["📊 Summary Reports", "🏆 Brand Reports", "🏷️ SKU Reports", "👥 Customer Channel", "👤 Sales Rep", "📽️ PowerPoint"])
            
            with tab1:
                st.markdown("### Sales Summary Report")
                st.markdown("Complete overview with MTD/YTD metrics, charts, and year-over-year comparisons")
                st.download_button(
                    label="📥 Download Sales Summary Report",
                    data=excel_output,
                    file_name=f"Sales_Summary_Report_{month_name}_{target_year}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True
                )
            
            with tab2:
                st.markdown("### Top 10 Brand Performance")
                col_b1, col_b2 = st.columns(2)
                
                with col_b1:
                    st.markdown("#### 📈 MTD Brands")
                    if brand_mtd_output:
                        st.download_button(
                            label=f"📥 Download MTD Report ({len(brand_mtd_df)} brands)",
                            data=brand_mtd_output,
                            file_name=f"Top10_Brand_MTD_Report_{month_name}_{target_year}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                        st.caption(f"✓ Top {len(brand_mtd_df)} brands by MTD sales")
                    else:
                        st.error("❌ Could not generate MTD brand report")
                
                with col_b2:
                    st.markdown("#### 📊 YTD Brands")
                    if brand_ytd_output:
                        st.download_button(
                            label=f"📥 Download YTD Report ({len(brand_ytd_df)} brands)",
                            data=brand_ytd_output,
                            file_name=f"Top10_Brand_YTD_Report_{month_name}_{target_year}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                        st.caption(f"✓ Top {len(brand_ytd_df)} brands by YTD sales")
                    else:
                        st.error("❌ Could not generate YTD brand report")
            
            with tab3:
                st.markdown("### Top 20 SKU Performance")
                col_s1, col_s2 = st.columns(2)
                
                with col_s1:
                    st.markdown("#### 📈 MTD SKUs")
                    if sku_mtd_output:
                        st.download_button(
                            label=f"📥 Download MTD Report ({len(sku_mtd_df)} SKUs)",
                            data=sku_mtd_output,
                            file_name=f"Top20_SKU_MTD_Report_{month_name}_{target_year}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                        st.caption(f"✓ Top {len(sku_mtd_df)} SKUs by MTD sales")
                    else:
                        st.error("❌ Could not generate MTD SKU report")
                
                with col_s2:
                    st.markdown("#### 📊 YTD SKUs")
                    if sku_ytd_output:
                        st.download_button(
                            label=f"📥 Download YTD Report ({len(sku_ytd_df)} SKUs)",
                            data=sku_ytd_output,
                            file_name=f"Top20_SKU_YTD_Report_{month_name}_{target_year}.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
                        st.caption(f"✓ Top {len(sku_ytd_df)} SKUs by YTD sales")
                    else:
                        st.error("❌ Could not generate YTD SKU report")
            
            with tab4:
                st.markdown("### 👥 Customer Channel Report")
                st.markdown("Top 10 customers by channel with MTD/YTD analysis")
                
                if channel_excel_output:
                    st.download_button(
                        label="📥 Download Customer Channel Report",
                        data=channel_excel_output,
                        file_name=f"Top10_Customers_by_Channel_{month_name}_{target_year}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )
                    st.caption("✓ Includes channel summary and Top 10 MTD/YTD customers per channel")
                else:
                    st.error("❌ Could not generate channel customer report")
            
            with tab5:
                st.markdown("### � Sales Rep Performance Report")
                st.markdown("MTD/YTD performance analysis by Sales Rep with charts")
                
                if sales_rep_excel_output:
                    st.download_button(
                        label="📥 Download Sales Rep Report",
                        data=sales_rep_excel_output,
                        file_name=f"Sales_Rep_Performance_{month_name}_{target_year}.xlsx",
                        mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                        use_container_width=True
                    )
                    st.caption("✓ Includes sales rep summary with MTD/YTD metrics and performance charts")
                else:
                    st.error("❌ Could not generate sales rep report")
            
            with tab6:
                st.markdown("### 📽️ PowerPoint Presentation")
                st.markdown("""
                **Comprehensive slide deck including:**
                - Executive summary with key metrics
                - MTD vs YTD performance charts
                - Top 10 brands analysis (MTD & YTD)
                - Top 20 SKUs breakdown (MTD & YTD)
                - Channel sales performance summary
                - Sales rep performance summary
                - Professional formatting and visualizations
                """)
                
                if ppt_output:
                    st.download_button(
                        label="📥 Download PowerPoint Presentation",
                        data=ppt_output,
                        file_name=f"Sales_Performance_Report_{month_name}_{target_year}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    st.success("✅ PowerPoint presentation ready for download!")
                    st.caption("Perfect for executive presentations and stakeholder meetings")
                else:
                    st.error("❌ Could not generate PowerPoint presentation")
            
    except Exception as e:
        st.error(f"❌ Error processing file: {str(e)}")
        st.exception(e)
else:
    st.info("👆 Please upload all required Excel files to get started")
    
    # Show sample information
    with st.expander("ℹ️ How to use this app"):
        st.markdown("""
        ### Instructions:
        1. **Upload your Sales Summary file** (.xls or .xlsx format)
        2. **Upload your ItemMaster file** (.xls or .xlsx format) - **Required**
           - This file should have "ItemId" and "Brand" columns
           - This is required for accurate brand identification in Top 10 Brands reports
        3. **Upload Sales Details file** (.xls or .xlsx format) - **Required**
           - Contains transaction-level sales data
        4. **Upload Sales Customer List file** (.xls or .xlsx format) - **Required**
           - Contains customer information with Channel column
        5. **Configure settings** in the sidebar:
           - Select target month for MTD calculations
           - Select target year for current data
           - Select comparison year for year-over-year analysis
        6. **View the results** in the dashboard
        7. **Download** the generated reports including:
           - Formatted summary table
           - Dashboard chart (4-panel visualization)
           - Top 10 Brands reports (MTD & YTD)
           - Top 20 SKUs reports (MTD & YTD)
           - Top 10 Customers by Channel (MTD & YTD)
           - Sales Rep Performance (MTD & YTD)
           - PowerPoint presentation
        
        ### Why is ItemMaster File Required?
        Different brands may share the same first 3 letters in their item numbers, 
        making it unreliable to identify brands by item codes alone. The ItemMaster 
        file provides accurate brand-to-item mappings for precise reporting.
        
        ### Expected Sales Summary File Format:
        - Header row should be at row 9 (0-indexed)
        - Should contain columns: Year, Period, Sales Amount, Cost of Sales
        - Items identified by alphanumeric codes
        - "Item Total:" rows separate different items
        
        ### Customer Channel Analysis:
        The **Sales Details** and **Sales Customer List** files are used to generate:
        - Channel-level MTD/YTD summary with charts
        - Top 10 customers per channel for MTD and YTD
        - YoY growth analysis by channel
        
        ### Sales Rep Performance:
        The **Sales Details** and **Sales Customer List** files are also used to generate:
        - Sales Rep MTD/YTD performance summary
        - Performance charts comparing current vs previous year
        - % Achieved and YoY Growth metrics by sales rep
        """)

# Footer
st.markdown("---")
st.markdown("Built with Streamlit 🎈 | Sales Report Generator v1.0")
